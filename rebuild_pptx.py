"""Rebuild PowerPoint slides 4-9 with proper formatting."""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import pandas as pd
import numpy as np

PPTX_PATH = r"d:\Context Engineering Pressing Analyst.pptx"
prs = Presentation(PPTX_PATH)

# Slide dimensions
SW = prs.slide_width   # 18288000 EMU = 20 inches
SH = prs.slide_height  # 10287000 EMU = 11.25 inches

# ── Colour palette ──────────────────────────────────────────────
WHITE  = RGBColor(255, 255, 255)
GREY   = RGBColor(170, 170, 170)
GREEN  = RGBColor(0, 200, 120)
DARK   = RGBColor(30, 30, 30)

FONT_NAME = "Calibri"

# ── Helper: add a text box ──────────────────────────────────────
def add_textbox(slide, left, top, width, height, text, font_size=13,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_NAME):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_rich_textbox(slide, left, top, width, height, runs_list,
                     font_name=FONT_NAME, alignment=PP_ALIGN.LEFT, spacing=Pt(4)):
    """
    runs_list: list of paragraphs, each paragraph is a list of (text, size, color, bold) tuples.
    """
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for pi, para_runs in enumerate(runs_list):
        if pi == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_after = spacing
        for ri, (text, size, color, bold) in enumerate(para_runs):
            run = p.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = font_name
            run.font.bold = bold
    return txBox


# ── Compute z-scores and wordalisation texts ────────────────────
df = pd.read_csv("data/pressing/pressing_detailed_metrics.csv")

METRIC_LABELS = {
    "chains_pm": "Pressing Chains Per Match",
    "recovery_opp_half_pm": "Regains In Opp. Half Per Match",
    "force_long_ball_pm": "Opp. Long Balls Forced Per Match",
    "press_break_rate_opp_half": "Press Break Rate In Opp. Half",
    "lead_to_goal_after_broken_pm": "Opp. Goals After Broken Press Per Match",
    "pass_accuracy_under_pressure": "Opp. Pass Accuracy Under Pressure",
    "line_breaking_pass_rate_under_pressure": "Opp. Line-Breaking Pass Rate Under Pressure",
    "ppda": "PPDA",
}

METRIC_CONSEQUENCES = {
    "chains_pm": {
        "outstanding": "The team starts pressing sequences in the opposition half far more than any other side.",
        "excellent": "The team starts pressing sequences in the opposition half significantly more often than most.",
        "good": "The team starts pressing sequences in the opposition half more often than average.",
        "average": "The team starts pressing sequences in the opposition half at a typical rate.",
        "below average": "The team starts pressing sequences in the opposition half less often than most.",
        "poor": "The team rarely starts pressing sequences in the opposition half.",
    },
    "recovery_opp_half_pm": {
        "outstanding": "The team wins the ball back in the opposition half far more than any other side.",
        "excellent": "The team wins the ball back in the opposition half significantly more often than most.",
        "good": "The team wins the ball back in the opposition half more often than average.",
        "average": "The team wins the ball back in the opposition half at a typical rate.",
        "below average": "The team wins the ball back in the opposition half less often than most.",
        "poor": "The team rarely wins the ball back in the opposition half.",
    },
    "force_long_ball_pm": {
        "outstanding": "The press forces opponents into long balls far more than any other side.",
        "excellent": "The press forces opponents into long balls significantly more often than most.",
        "good": "The press forces opponents into long balls more often than average.",
        "average": "The press forces opponents into long balls at a typical rate.",
        "below average": "The press forces opponents into long balls less often than most.",
        "poor": "The press rarely forces opponents into long balls.",
    },
    "press_break_rate_opp_half": {
        "outstanding": "Opponents almost never escape the press in the opposition half.",
        "excellent": "Opponents rarely escape the press in the opposition half.",
        "good": "Opponents escape the press in the opposition half less often than against most teams.",
        "average": "Opponents escape the press in the opposition half at a typical rate.",
        "below average": "Opponents escape the press in the opposition half more often than against most teams.",
        "poor": "Opponents escape the press in the opposition half frequently - the press is easily broken.",
    },
    "lead_to_goal_after_broken_pm": {
        "outstanding": "When the press is broken, opponents almost never score - the rest-defence absorbs danger completely.",
        "excellent": "When the press is broken, opponents rarely score - the rest-defence contains the damage well.",
        "good": "When the press is broken, opponents score less often than against most teams.",
        "average": "When the press is broken, opponents score at a typical rate.",
        "below average": "When the press is broken, opponents score more often than against most teams.",
        "poor": "When the press is broken, opponents frequently score - the rest-defence is exposed.",
    },
    "pass_accuracy_under_pressure": {
        "outstanding": "Opponents struggle badly to complete passes under this team pressure - passing accuracy collapses.",
        "excellent": "Opponents complete far fewer passes under this team pressure - passing accuracy drops sharply.",
        "good": "Opponents find it harder to complete passes under this team pressure than against most sides.",
        "average": "Opponents complete passes under this team pressure at a typical rate.",
        "below average": "Opponents still complete passes fairly comfortably under this team pressure.",
        "poor": "Opponents pass comfortably under this team pressure - the press does little to disrupt passing accuracy.",
    },
    "line_breaking_pass_rate_under_pressure": {
        "outstanding": "Opponents almost never play through the defensive lines under this team pressure.",
        "excellent": "Opponents rarely play through the defensive lines under this team pressure - the pressing structure holds its shape.",
        "good": "Opponents break through the defensive lines under this team pressure less often than against most sides.",
        "average": "Opponents break through the defensive lines under this team pressure at a typical rate.",
        "below average": "Opponents break through the defensive lines under this team pressure more often than against most sides.",
        "poor": "Opponents play through the defensive lines comfortably under this team pressure.",
    },
    "ppda": {
        "outstanding": "The team steps in extremely early, allowing opponents almost no passes before a defensive action.",
        "excellent": "The team steps in early, allowing opponents very few passes before a defensive action.",
        "good": "The team steps in to break opponent sequences quicker than most.",
        "average": "The team allows a typical number of opponent passes before stepping in with a defensive action.",
        "below average": "The team is slow to step in, allowing opponents more passes before a defensive action than most.",
        "poor": "The team is very slow to step in - opponents circulate freely before any defensive action arrives.",
    },
}

# Inverted metrics (lower raw = better)
INVERTED = {
    "ppda", "press_break_rate_opp_half", "lead_to_goal_after_broken_pm",
    "pass_accuracy_under_pressure", "line_breaking_pass_rate_under_pressure",
}

metrics = list(METRIC_LABELS.keys())

def describe_level(z):
    thresholds = [1.5, 1, 0.5, -0.5, -1]
    words = ["outstanding", "excellent", "good", "average", "below average", "poor"]
    i = 0
    while i < len(thresholds) and z < thresholds[i]:
        i += 1
    return words[i]

def team_wordalisation(team_name):
    row = df[df["team"] == team_name].iloc[0]
    lines = []
    standouts = []
    concerns = []
    for m in metrics:
        raw = row[m]
        col_mean = df[m].mean()
        col_std = df[m].std()
        z = (raw - col_mean) / col_std
        if m in INVERTED:
            z = -z
        level = describe_level(z)
        label = METRIC_LABELS[m]
        consequence = METRIC_CONSEQUENCES[m][level]
        lines.append((consequence, label))
        if z > 1.0:
            standouts.append(label)
        elif z < -0.5:
            concerns.append(label)
    return lines, standouts, concerns


# ── Generate wordalisation for 3 teams ──────────────────────────
fulham_lines, fulham_str, fulham_con = team_wordalisation("Fulham")
bmouth_lines, bmouth_str, bmouth_con = team_wordalisation("Bournemouth")
ipswich_lines, ipswich_str, ipswich_con = team_wordalisation("Ipswich Town")

# Pre-written coaching briefings (from few-shot quality outputs)
FULHAM_BRIEFING = (
    "Fulham press at a standard rate without standing out in any particular area "
    "of high-press intensity or frequency. Their clear strength is behind the press: "
    "when opponents do escape, they rarely convert into goals, suggesting a well-organised "
    "rest-defence that absorbs danger effectively. However, the press itself offers limited "
    "disruption - opponents escape at a higher-than-average rate, pass comfortably under "
    "pressure, and are not forced into long balls with any regularity. Overall, Fulham "
    "prioritise defensive resilience and damage limitation over aggressive high pressing, "
    "accepting that opponents will build through the initial engagement but backing themselves "
    "to contain the consequences."
)

BOURNEMOUTH_BRIEFING = (
    "Bournemouth press with real intensity, starting sequences high and often, stepping in "
    "early to break opponent build-up and winning the ball back in the opposition half at one "
    "of the highest rates in the league. The press locks opponents in effectively, forces long "
    "balls, and when the press is broken the rest-defence limits the damage - a well-rounded, "
    "aggressive system with genuine bite. Opponents can still play through the lines at a fairly "
    "typical rate, meaning the press disrupts passing but does not fully shut down vertical "
    "progression. Overall, Bournemouth commit heavily to high pressing and back it with solid "
    "cover, trading some vulnerability to line-breaking passes for dominant territorial pressure."
)

IPSWICH_BRIEFING = (
    "Ipswich press at a standard rate but the structure behind that press is fragile - the press "
    "is very slow to step in, allowing opponents to circulate freely, and rarely forces them into "
    "long balls. Opponents escape the press more often than against most sides and retain their "
    "passing quality comfortably, meaning the engagement offers minimal genuine disruption. The "
    "one redeeming element is that when opponents do escape, they score less often than against "
    "most teams, suggesting some capacity to absorb danger deeper. Overall, Ipswich run a press "
    "that lacks structural intensity and coordination, relying on deeper defensive resilience "
    "rather than front-foot disruption."
)


# ── CLEAN UP: Remove my previously added text boxes on slides 4-9 ─────
for slide_idx in [3, 4, 5, 6, 8]:  # 0-indexed: slides 4,5,6,7,9
    slide = prs.slides[slide_idx]
    to_remove = []
    for shape in slide.shapes:
        if shape.name.startswith("TextBox"):
            to_remove.append(shape)
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

# Also clean slide 7 (index 6) text boxes
slide7 = prs.slides[6]
to_remove = []
for shape in slide7.shapes:
    if shape.name.startswith("TextBox"):
        to_remove.append(shape)
for shape in to_remove:
    sp = shape._element
    sp.getparent().remove(sp)


# ── SLIDES 4-6: Reposition images + add wordalisation text ─────

def build_team_slide(slide, team_name, briefing, lines, standouts, concerns, image_name):
    """Reposition image to right, add briefing text on left."""

    # Find the image and reposition it to right side
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            # Right half: starts at 52% of width, takes ~46%
            shape.left = int(SW * 0.50)
            shape.top  = int(SH * 0.15)
            shape.width = int(SW * 0.48)
            shape.height = int(SH * 0.80)

    # LEFT SIDE: Briefing text
    left_margin = int(SW * 0.025)
    text_width  = int(SW * 0.46)
    top_start   = int(SH * 0.17)

    # Build paragraphs for the rich textbox
    paragraphs = []

    # Briefing text (main body)
    paragraphs.append([
        (briefing, 14, WHITE, False)
    ])

    # Spacer
    paragraphs.append([(" ", 8, WHITE, False)])

    # Synthesized input label
    paragraphs.append([
        ("What the LLM receives:", 12, GREEN, True)
    ])

    # Show first 4 consequence lines as sample
    for consequence, label in lines[:4]:
        para = [
            (consequence, 11, GREY, False),
            ("  (" + label + ")", 10, RGBColor(120, 120, 120), False),
        ]
        paragraphs.append(para)

    # "... and N more metrics"
    remaining = len(lines) - 4
    if remaining > 0:
        paragraphs.append([
            ("... and {} more metrics".format(remaining), 10, RGBColor(100, 100, 100), True)
        ])

    # Key strengths / concerns
    if standouts:
        paragraphs.append([(" ", 6, WHITE, False)])
        paragraphs.append([
            ("Key strengths: ", 11, GREEN, True),
            (", ".join(standouts), 11, WHITE, False),
        ])
    if concerns:
        paragraphs.append([
            ("Key concerns: ", 11, RGBColor(255, 100, 100), True),
            (", ".join(concerns), 11, WHITE, False),
        ])

    add_rich_textbox(slide, left_margin, top_start, text_width, int(SH * 0.78),
                     paragraphs, alignment=PP_ALIGN.LEFT)


# Slide 4 - Fulham
build_team_slide(prs.slides[3], "Fulham", FULHAM_BRIEFING,
                 fulham_lines, fulham_str, fulham_con, None)

# Slide 5 - Bournemouth
build_team_slide(prs.slides[4], "Bournemouth", BOURNEMOUTH_BRIEFING,
                 bmouth_lines, bmouth_str, bmouth_con, None)

# Slide 6 - Ipswich Town
build_team_slide(prs.slides[5], "Ipswich Town", IPSWICH_BRIEFING,
                 ipswich_lines, ipswich_str, ipswich_con, None)


# ── SLIDE 7: How the data reaches the LLM ──────────────────────
slide7 = prs.slides[6]

# Title
add_textbox(slide7, int(SW * 0.025), int(SH * 0.02), int(SW * 0.95), int(SH * 0.08),
            "How the Data Reaches the LLM", font_size=28, color=WHITE, bold=True)

# Subtitle
add_textbox(slide7, int(SW * 0.025), int(SH * 0.10), int(SW * 0.95), int(SH * 0.06),
            "Each metric is converted into a direct consequence statement - no room for misinterpretation",
            font_size=14, color=GREY)

# LEFT COLUMN: "Before" (ambiguous)
col_left = int(SW * 0.025)
col_width = int(SW * 0.44)

add_textbox(slide7, col_left, int(SH * 0.20), col_width, int(SH * 0.05),
            "BEFORE: Ambiguous level descriptions", font_size=14, color=RGBColor(255, 100, 100), bold=True)

before_lines = [
    [("The team was ", 12, GREY, False), ("poor", 12, RGBColor(255, 100, 100), True),
     (" in Opp. Pass Accuracy Under Pressure", 12, GREY, False)],
    [("", 6, GREY, False)],
    [("Problem: ", 12, RGBColor(255, 100, 100), True),
     ("Does 'poor' mean opponents pass well or badly?", 12, WHITE, False)],
    [("The LLM has to guess direction - and often gets it wrong.", 12, WHITE, False)],
]

add_rich_textbox(slide7, col_left, int(SH * 0.27), col_width, int(SH * 0.25),
                 before_lines)

# RIGHT COLUMN: "After" (consequence statements)
col_right = int(SW * 0.52)

add_textbox(slide7, col_right, int(SH * 0.20), col_width, int(SH * 0.05),
            "AFTER: Direct consequence statements", font_size=14, color=GREEN, bold=True)

after_lines = [
    [("Opponents pass comfortably under this team's pressure - the press does little to disrupt passing accuracy.", 12, WHITE, False)],
    [("", 6, GREY, False)],
    [("Benefit: ", 12, GREEN, True),
     ("No ambiguity. The LLM reads plain football language.", 12, WHITE, False)],
    [("48 pre-written statements (8 metrics x 6 levels) cover every case.", 12, WHITE, False)],
]

add_rich_textbox(slide7, col_right, int(SH * 0.27), col_width, int(SH * 0.25),
                 after_lines)

# BOTTOM: Full Arsenal example
add_textbox(slide7, int(SW * 0.025), int(SH * 0.56), int(SW * 0.95), int(SH * 0.05),
            "Full example input for Arsenal:", font_size=14, color=GREEN, bold=True)

arsenal_lines, arsenal_str, arsenal_con = team_wordalisation("Arsenal")
arsenal_paragraphs = []
for consequence, label in arsenal_lines:
    arsenal_paragraphs.append([
        (consequence, 10, GREY, False),
        ("  (" + label + ")", 9, RGBColor(100, 100, 100), False),
    ])

if arsenal_str:
    arsenal_paragraphs.append([
        ("Key strengths: ", 10, GREEN, True),
        (", ".join(arsenal_str), 10, WHITE, False),
    ])
if arsenal_con:
    arsenal_paragraphs.append([
        ("Key concerns: ", 10, RGBColor(255, 100, 100), True),
        (", ".join(arsenal_con), 10, WHITE, False),
    ])

add_rich_textbox(slide7, int(SW * 0.025), int(SH * 0.62), int(SW * 0.95), int(SH * 0.36),
                 arsenal_paragraphs)


# ── SLIDE 9: The Prompt ─────────────────────────────────────────
slide9 = prs.slides[8]

# Remove the empty placeholders too
to_remove = []
for shape in slide9.shapes:
    if hasattr(shape, "text") and shape.text.strip() == "":
        to_remove.append(shape)
for shape in to_remove:
    sp = shape._element
    sp.getparent().remove(sp)

# Title
add_textbox(slide9, int(SW * 0.025), int(SH * 0.02), int(SW * 0.95), int(SH * 0.08),
            "The Prompt", font_size=28, color=WHITE, bold=True)

# System message
sys_msg = (
    "You are a tactical data analyst embedded in a Premier League coaching staff. "
    "You write concise pressing briefings for the head coach. Your language is direct "
    "and action-oriented. You write in British English and refer to the sport as football."
)

# User prompt
user_prompt = (
    "Please use the pressing profile enclosed with ``` to give a concise 4-sentence briefing "
    "of this team's pressing style, strengths, and weaknesses. "
    "The first sentence should describe the team's pressing identity and how it is executed. "
    "The second sentence should describe pressing strengths. "
    "The third sentence should describe pressing limitations. "
    "The fourth sentence should summarise what this pressing approach prioritises and trades off. "
    "Do not invent consequences not in the data (e.g. rapid restarts, transition speed, possession recycling). "
    "Do not include metric names, level labels (e.g. excellent, poor), or parenthetical references in your output. "
    "Write as a tactical analyst would in a briefing to coaching staff."
)

prompt_paragraphs = [
    # System message header
    [("System Message", 16, GREEN, True)],
    [(sys_msg, 12, WHITE, False)],
    # Spacer
    [(" ", 8, WHITE, False)],
    # User prompt header
    [("User Prompt", 16, GREEN, True)],
    [(user_prompt, 12, WHITE, False)],
    # Spacer
    [(" ", 8, WHITE, False)],
    # Few-shot header
    [("Few-Shot Examples", 16, GREEN, True)],
    [("3 examples (Liverpool, Arsenal, Nottingham) teach the LLM:", 12, WHITE, False)],
    [("  - Correct calibration of language to metric levels", 11, GREY, False)],
    [("  - 4-sentence structure and coaching-staff tone", 11, GREY, False)],
    [("  - How to synthesise across all 8 metrics without inventing", 11, GREY, False)],
    # Spacer
    [(" ", 8, WHITE, False)],
    # Key design decisions
    [("Key Design Decisions", 16, GREEN, True)],
    [("  - Direct consequence statements eliminate metric-direction confusion", 11, WHITE, False)],
    [("  - Explicit prohibition list prevents hallucinated concepts", 11, WHITE, False)],
    [("  - Metric names in parentheses for traceability, banned from output", 11, WHITE, False)],
    [("  - Standouts and concerns guide the LLM to emphasise the right metrics", 11, WHITE, False)],
]

add_rich_textbox(slide9, int(SW * 0.025), int(SH * 0.12), int(SW * 0.95), int(SH * 0.85),
                 prompt_paragraphs)


# ── Save ────────────────────────────────────────────────────────
OUT_PATH = r"d:\Context Engineering Pressing Analyst v2.pptx"
prs.save(OUT_PATH)
print("Done! Saved to", OUT_PATH)
