"""
Build Pressing Analyst presentation — styled after the Physical Analyst deck.
Same narrative arc: examples first, then step-by-step, then pipeline.
Run:  python build_pressing_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd
from scipy.stats import zscore

# ── Palette ─────────────────────────────────────────────────────
BG       = RGBColor(20, 20, 38)
CARD     = RGBColor(32, 32, 56)
ACCENT   = RGBColor(0, 200, 120)
ACCENT2  = RGBColor(0, 160, 96)
RED      = RGBColor(255, 90, 90)
AMBER    = RGBColor(255, 190, 60)
WHITE    = RGBColor(255, 255, 255)
BODY     = RGBColor(210, 210, 225)
MUTED    = RGBColor(140, 140, 170)
DIM      = RGBColor(90, 90, 120)
FONT     = "Calibri"
MONO     = "Consolas"

# ── Dimensions ──────────────────────────────────────────────────
SW = Inches(13.333)
SH = Inches(7.5)
M  = Inches(0.6)
CW = SW - 2 * M

# ── Data ────────────────────────────────────────────────────────
df = pd.read_csv("data/pressing/pressing_detailed_metrics.csv")
df_league = pd.read_csv("data/pressing/pressing_league_table.csv")
N = len(df)

METRICS = [
    "chains_pm", "recovery_opp_half_pm", "force_long_ball_pm",
    "press_break_rate_opp_half", "lead_to_goal_after_broken_pm",
    "pass_accuracy_under_pressure", "line_breaking_pass_rate_under_pressure", "ppda",
]
INVERTED = {
    "ppda", "press_break_rate_opp_half", "lead_to_goal_after_broken_pm",
    "pass_accuracy_under_pressure", "line_breaking_pass_rate_under_pressure",
}
METRIC_LABELS = {
    "chains_pm": "Pressing Chains Per Match",
    "recovery_opp_half_pm": "Regains In Opp. Half Per Match",
    "force_long_ball_pm": "Opp. Long Balls Forced Per Match",
    "press_break_rate_opp_half": "Press Break Rate In Opp. Half",
    "lead_to_goal_after_broken_pm": "Opp. Goals After Broken Press Per Match",
    "pass_accuracy_under_pressure": "Opp. Pass Accuracy Under Pressure",
    "line_breaking_pass_rate_under_pressure": "Opp. Line-Breaking Pass Rate Under Pressure",
    "ppda": "PPDA",
}
SHORT_LABELS = {
    "chains_pm": "Chains/M",
    "recovery_opp_half_pm": "Regains",
    "force_long_ball_pm": "Long Balls",
    "press_break_rate_opp_half": "Break Rate",
    "lead_to_goal_after_broken_pm": "Goals After Break",
    "pass_accuracy_under_pressure": "Pass Acc.",
    "line_breaking_pass_rate_under_pressure": "Line Break",
    "ppda": "PPDA",
}
METRIC_CONSEQUENCES = {
    "chains_pm": {
        "outstanding": "The team starts pressing sequences in the opposition half far more than any other side.",
        "excellent": "The team starts pressing sequences in the opposition half significantly more often than most.",
        "good": "The team starts pressing sequences in the opposition half more often than average.",
        "average": "The team starts pressing sequences in the opposition half at a typical rate.",
        "below average": "The team starts pressing sequences in the opposition half less often than most.",
        "poor": "The team rarely starts pressing sequences in the opposition half.",
    },
    "recovery_opp_half_pm": {
        "outstanding": "The team wins the ball back in the opposition half far more than any other side.",
        "excellent": "The team wins the ball back in the opposition half significantly more often than most.",
        "good": "The team wins the ball back in the opposition half more often than average.",
        "average": "The team wins the ball back in the opposition half at a typical rate.",
        "below average": "The team wins the ball back in the opposition half less often than most.",
        "poor": "The team rarely wins the ball back in the opposition half.",
    },
    "force_long_ball_pm": {
        "outstanding": "The press forces opponents into long balls far more than any other side.",
        "excellent": "The press forces opponents into long balls significantly more often than most.",
        "good": "The press forces opponents into long balls more often than average.",
        "average": "The press forces opponents into long balls at a typical rate.",
        "below average": "The press forces opponents into long balls less often than most.",
        "poor": "The press rarely forces opponents into long balls.",
    },
    "press_break_rate_opp_half": {
        "outstanding": "Opponents almost never escape the press in the opposition half.",
        "excellent": "Opponents rarely escape the press in the opposition half.",
        "good": "Opponents escape the press in the opposition half less often than against most teams.",
        "average": "Opponents escape the press in the opposition half at a typical rate.",
        "below average": "Opponents escape the press in the opposition half more often than against most teams.",
        "poor": "Opponents escape the press in the opposition half frequently.",
    },
    "lead_to_goal_after_broken_pm": {
        "outstanding": "When the press is broken, opponents almost never score.",
        "excellent": "When the press is broken, opponents rarely score.",
        "good": "When the press is broken, opponents score less often than against most teams.",
        "average": "When the press is broken, opponents score at a typical rate.",
        "below average": "When the press is broken, opponents score more often than against most teams.",
        "poor": "When the press is broken, opponents frequently score.",
    },
    "pass_accuracy_under_pressure": {
        "outstanding": "Opponents struggle badly to complete passes under this team\u2019s pressure.",
        "excellent": "Opponents complete far fewer passes under this team\u2019s pressure.",
        "good": "Opponents find it harder to complete passes under this team\u2019s pressure than against most sides.",
        "average": "Opponents complete passes under this team\u2019s pressure at a typical rate.",
        "below average": "Opponents still complete passes fairly comfortably under this team\u2019s pressure.",
        "poor": "Opponents pass comfortably under this team\u2019s pressure.",
    },
    "line_breaking_pass_rate_under_pressure": {
        "outstanding": "Opponents almost never play through the defensive lines under this team\u2019s pressure.",
        "excellent": "Opponents rarely play through the defensive lines under this team\u2019s pressure.",
        "good": "Opponents break through the defensive lines less often than against most sides.",
        "average": "Opponents break through the defensive lines at a typical rate.",
        "below average": "Opponents break through the defensive lines more often than against most sides.",
        "poor": "Opponents play through the defensive lines comfortably.",
    },
    "ppda": {
        "outstanding": "The team steps in extremely early, allowing opponents almost no passes before a defensive action.",
        "excellent": "The team steps in early, allowing opponents very few passes before a defensive action.",
        "good": "The team steps in to break opponent sequences quicker than most.",
        "average": "The team allows a typical number of opponent passes before stepping in.",
        "below average": "The team is slow to step in, allowing opponents more passes than most.",
        "poor": "The team is very slow to step in \u2014 opponents circulate freely.",
    },
}

# ── Helpers ─────────────────────────────────────────────────────
def describe_level(z):
    for thr, w in zip([1.5, 1, 0.5, -0.5, -1],
                      ["Outstanding", "Excellent", "Good", "Average", "Below average"]):
        if z >= thr:
            return w
    return "Poor"

def level_color(lev):
    return {"Outstanding": ACCENT, "Excellent": ACCENT, "Good": RGBColor(100, 210, 160),
            "Average": BODY, "Below average": AMBER, "Poor": RED}.get(lev, BODY)

def ordinal(n):
    if 11 <= (n % 100) <= 13: return f"{n}th"
    return f"{n}{({1:'st',2:'nd',3:'rd'}).get(n%10,'th')}"

def fmt(m, v):
    v = float(v)
    if m in {"press_break_rate_opp_half","pass_accuracy_under_pressure","line_breaking_pass_rate_under_pressure"}:
        return f"{v*100:.1f}%"
    return f"{v:.2f}"

def team_profile(name):
    row = df[df["team"] == name].iloc[0]
    zdf = df[METRICS].apply(zscore, nan_policy="omit")
    for m in INVERTED: zdf[m] *= -1
    idx = df[df["team"] == name].index[0]
    results = []
    for m in METRICS:
        z = float(zdf.loc[idx, m])
        lev = describe_level(z)
        lab = METRIC_LABELS[m]
        short = SHORT_LABELS[m]
        cons = METRIC_CONSEQUENCES[m][lev.lower()]
        rank = int(df[m].rank(ascending=(m in INVERTED)).loc[idx])
        results.append((m, lab, short, fmt(m, row[m]), rank, lev, cons, z))
    bp, br = None, None
    if "Block %" in df_league.columns:
        lr = df_league[df_league["Team"] == name]
        if not lr.empty:
            bp = float(lr["Block %"].iloc[0])
            br = int((df_league["Block %"] > bp).sum() + 1)
    return results, bp, br

# ── Shape helpers ───────────────────────────────────────────────
def _noborder(s): s.line.fill.background()

def rect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill; _noborder(s)
    if s.adjustments: s.adjustments[0] = 0.02
    return s

def hline(slide, l, t, w, c=ACCENT, th=Pt(2)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, th)
    s.fill.solid(); s.fill.fore_color.rgb = c; _noborder(s)

def setup(slide, num, total=13):
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = BG
    hline(slide, Emu(0), Emu(0), SW, ACCENT, Pt(4))
    tb(slide, Inches(0.4), SH-Inches(0.42), Inches(4), Inches(0.3),
       "Context Engineering for Football  \u00b7  Twelve Football  \u00b7  2026", sz=8, c=DIM)
    tb(slide, SW-Inches(1), SH-Inches(0.42), Inches(0.7), Inches(0.3),
       f"{num}/{total}", sz=8, c=DIM, al=PP_ALIGN.RIGHT)

def tb(slide, l, t, w, h, text, sz=13, c=WHITE, b=False, al=PP_ALIGN.LEFT, f=FONT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = c; p.font.name = f; p.font.bold = b; p.alignment = al
    return box

def rich(slide, l, t, w, h, paras, al=PP_ALIGN.LEFT, sp=Pt(4)):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for pi, runs in enumerate(paras):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = al; p.space_after = sp
        for text, sz, clr, bld, *extra in runs:
            r = p.add_run(); r.text = text; r.font.size = Pt(sz)
            r.font.color.rgb = clr; r.font.name = extra[0] if extra else FONT; r.font.bold = bld
    return box

# ── Build ───────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SW; prs.slide_height = SH
blank = prs.slide_layouts[6]

BRIEFINGS = {
    "Bournemouth": (
        "Bournemouth press with real intensity, starting sequences high and often, stepping in "
        "early to break opponent build-up and winning the ball back in the opposition half at one "
        "of the highest rates in the league. The press locks opponents in effectively, forces long "
        "balls, and when the press is broken the rest-defence limits the damage \u2014 a well-rounded, "
        "aggressive system with genuine bite. Opponents can still play through the lines at a fairly "
        "typical rate, meaning the press disrupts passing but does not fully shut down vertical "
        "progression. Overall, Bournemouth commit heavily to high pressing and back it with solid "
        "cover, trading some vulnerability to line-breaking passes for dominant territorial pressure."
    ),
    "Fulham": (
        "Fulham press at a standard rate without standing out in any particular area "
        "of high-press intensity or frequency. Their clear strength is behind the press: "
        "when opponents do escape, they rarely convert into goals, suggesting a well-organised "
        "rest-defence that absorbs danger effectively. However, the press itself offers limited "
        "disruption \u2014 opponents escape at a higher-than-average rate, pass comfortably under "
        "pressure, and are not forced into long balls with any regularity. Overall, Fulham "
        "prioritise defensive resilience and damage limitation over aggressive high pressing, "
        "accepting that opponents will build through the initial engagement but backing themselves "
        "to contain the consequences."
    ),
    "Ipswich Town": (
        "Ipswich press at a standard rate but the structure behind that press is fragile \u2014 the press "
        "is very slow to step in, allowing opponents to circulate freely, and rarely forces them into "
        "long balls. Opponents escape the press more often than against most sides and retain their "
        "passing quality comfortably, meaning the engagement offers minimal genuine disruption. The "
        "one redeeming element is that when opponents do escape, they score less often than against "
        "most teams, suggesting some capacity to absorb danger deeper. Overall, Ipswich run a press "
        "that lacks structural intensity and coordination, relying on deeper defensive resilience "
        "rather than front-foot disruption."
    ),
}


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); setup(s, 1)
rect(s, Emu(0), Inches(2.2), SW, Inches(2.8), CARD)
hline(s, Inches(4.5), Inches(2.35), Inches(4.3), ACCENT, Pt(3))
tb(s, M, Inches(2.5), CW, Inches(0.7),
   "Pressing Analyst", sz=44, c=WHITE, b=True, al=PP_ALIGN.CENTER)
tb(s, M, Inches(3.3), CW, Inches(0.5),
   "Wordalization of Twelve Pressing Data", sz=18, c=MUTED, al=PP_ALIGN.CENTER)
tb(s, M, Inches(4.1), CW, Inches(0.4),
   "Pressing Volume  \u00b7  Disruption  \u00b7  Resilience  \u00b7  Engagement Speed",
   sz=14, c=ACCENT, al=PP_ALIGN.CENTER)
tb(s, M, Inches(5.5), CW, Inches(0.3),
   "Kaan G\u00fcler  \u00b7  Hugo Vicente  \u00b7  Sara Bentelli",
   sz=12, c=DIM, al=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDES 2-4 — TEAM EXAMPLES (same pattern as Physical slides 2-4)
# ════════════════════════════════════════════════════════════════
def build_example_slide(name, num):
    s = prs.slides.add_slide(blank); setup(s, num)
    tb(s, M, Inches(0.25), CW, Inches(0.5), name, sz=30, c=WHITE, b=True)
    tb(s, M, Inches(0.75), Inches(6), Inches(0.3),
       f"Premier League  \u00b7  2024/25", sz=12, c=MUTED)

    profile, bp, br = team_profile(name)

    # Left side: metric level badges
    y = Inches(1.3)
    for m, lab, short, val, rank, lev, cons, z in profile:
        clr = level_color(lev)
        rect(s, M, y, Inches(1.6), Inches(0.4), CARD)
        tb(s, Inches(0.7), y + Inches(0.05), Inches(1.5), Inches(0.3),
           short, sz=10, c=MUTED, b=True)
        rect(s, Inches(2.3), y, Inches(1.6), Inches(0.4), CARD)
        tb(s, Inches(2.35), y + Inches(0.05), Inches(1.5), Inches(0.3),
           lev, sz=10, c=clr, b=True)
        y += Inches(0.48)

    # Block % if available
    if bp is not None:
        rect(s, M, y + Inches(0.1), Inches(3.5), Inches(0.4), CARD)
        hline(s, M, y + Inches(0.1), Inches(3.5), ACCENT2, Pt(2))
        tb(s, Inches(0.7), y + Inches(0.15), Inches(3.2), Inches(0.3),
           f"Block %: {bp:.1f}% ({ordinal(br)}/{N})", sz=10, c=ACCENT, b=True)

    # Right side: LLM Output
    rect(s, Inches(4.5), Inches(1.3), Inches(8.2), Inches(5.5), CARD)
    hline(s, Inches(4.5), Inches(1.3), Inches(8.2), ACCENT, Pt(2))
    tb(s, Inches(4.7), Inches(1.4), Inches(3), Inches(0.3),
       "LLM Output", sz=12, c=ACCENT, b=True)
    tb(s, Inches(4.7), Inches(1.8), Inches(7.8), Inches(4.8),
       BRIEFINGS[name], sz=12, c=BODY)


build_example_slide("Bournemouth", 2)
build_example_slide("Fulham", 3)
build_example_slide("Ipswich Town", 4)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — STEP 1: TELL IT WHO IT IS (System Prompt)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); setup(s, 5)
tb(s, M, Inches(0.3), CW, Inches(0.5),
   "Step 1: Tell It Who It Is", sz=30, c=WHITE, b=True)
tb(s, M, Inches(0.85), CW, Inches(0.3),
   "The system prompt defines the persona", sz=14, c=MUTED)

rect(s, M, Inches(1.5), Inches(12), Inches(2.5), CARD)
hline(s, M, Inches(1.5), Inches(12), ACCENT, Pt(2))
tb(s, Inches(0.8), Inches(1.6), Inches(2), Inches(0.3),
   "System Prompt", sz=14, c=ACCENT, b=True)
sys_prompt = (
    '"You are a tactical data analyst embedded in a Premier League coaching staff.\n'
    'You write concise pressing briefings for the head coach.\n'
    'Your language is direct and action-oriented \u2014 you focus on what the data means tactically.\n'
    'You write in British English and refer to the sport as football."'
)
tb(s, Inches(0.8), Inches(2.0), Inches(11.2), Inches(1.8),
   sys_prompt, sz=12, c=BODY, f=MONO)

# Comparison callout
rect(s, M, Inches(4.4), Inches(12), Inches(2.3), CARD)
p = [
    [("What\u2019s different from Physical Analyst?", 14, WHITE, True)],
    [("", 4, WHITE, False)],
    [("\u2022  Physical: ", 12, MUTED, False), ("\"experienced physical performance analyst\"", 12, BODY, False, MONO),
     (" \u2014 profiles individual players", 12, MUTED, False)],
    [("\u2022  Pressing: ", 12, ACCENT, False), ("\"tactical data analyst on coaching staff\"", 12, BODY, False, MONO),
     (" \u2014 profiles entire teams", 12, ACCENT, False)],
    [("", 6, WHITE, False)],
    [("\u2022  Physical writes for recruitment departments; Pressing writes for the coaching room", 12, MUTED, False)],
    [("\u2022  Football/soccer grounding pair: ", 12, MUTED, False),
     ("\"Do you refer to the game as soccer or football?\" \u2192 \"Football.\"", 12, BODY, False, MONO)],
]
rich(s, Inches(0.8), Inches(4.55), Inches(11.5), Inches(2), p, sp=Pt(4))


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — STEP 2: TELL IT WHAT IT KNOWS (Q&A Pairs)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); setup(s, 6)
tb(s, M, Inches(0.3), CW, Inches(0.5),
   "Step 2: Tell It What It Knows", sz=30, c=WHITE, b=True)
tb(s, M, Inches(0.85), CW, Inches(0.3),
   "19 Q&A pairs teaching the LLM about pressing metrics \u2014 examples below", sz=14, c=MUTED)

qas = [
    ("Q:  What does PPDA measure?",
     "PPDA is opponent completed passes in the high block divided by this team\u2019s meaningful "
     "pressing actions in the same zone. Lower is more aggressive."),
    ("Q:  What does it mean if a team has low PPDA but high press break rate?",
     "This suggests a team that steps in early and aggressively but gets bypassed often. "
     "The engagement is frequent but structurally fragile."),
    ("Q:  What does high/medium block percentage tell us?",
     "It shows how much time a team spends defending in a high or medium block. "
     "This is a positional preference, not a quality judgement."),
]
y = Inches(1.5)
for q, a in qas:
    rect(s, M, y, Inches(12), Inches(1.5), CARD)
    tb(s, Inches(0.8), y + Inches(0.1), Inches(11), Inches(0.3),
       q, sz=12, c=ACCENT, b=True)
    tb(s, Inches(0.8), y + Inches(0.5), Inches(11), Inches(0.9),
       a, sz=11, c=BODY)
    y += Inches(1.65)

# Category footer
rect(s, M, Inches(6.5), Inches(12), Inches(0.55), CARD)
p = [
    [("Categories: ", 10, MUTED, True),
     ("metric definitions (8)  \u00b7  metric combinations (5)  \u00b7  interpretation rules (4)  \u00b7  context guidance (2)", 10, BODY, False)],
]
rich(s, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.4), p)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — STEP 3: TELL IT WHAT DATA TO USE (Consequence Sentences)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); setup(s, 7)
tb(s, M, Inches(0.3), CW, Inches(0.5),
   "Step 3: Tell It What Data to Use", sz=30, c=WHITE, b=True)
tb(s, M, Inches(0.85), CW, Inches(0.3),
   "Convert z-scores into consequence sentences \u2014 not just word labels", sz=14, c=MUTED)

# Left: Physical approach
rect(s, M, Inches(1.5), Inches(5.7), Inches(3.5), CARD)
hline(s, M, Inches(1.5), Inches(5.7), AMBER, Pt(2))
p = [
    [("Physical Analyst approach", 14, AMBER, True)],
    [("", 4, WHITE, False)],
    [("He was ", 12, BODY, False), ("excellent", 12, AMBER, True),
     (" in Speed.", 12, BODY, False)],
    [("He was ", 12, BODY, False), ("poor", 12, RED, True),
     (" in Agility.", 12, BODY, False)],
    [("", 8, WHITE, False)],
    [("Word label only.", 12, MUTED, False)],
    [("LLM must infer what \u2018poor\u2019 means in context.", 12, MUTED, False)],
]
rich(s, Inches(0.8), Inches(1.65), Inches(5.3), Inches(3.2), p, sp=Pt(4))

# Right: Our approach
rect(s, Inches(6.8), Inches(1.5), Inches(5.9), Inches(3.5), CARD)
hline(s, Inches(6.8), Inches(1.5), Inches(5.9), ACCENT, Pt(2))
p = [
    [("Pressing Analyst approach", 14, ACCENT, True)],
    [("", 4, WHITE, False)],
    [("The team steps in early, allowing opponents", 12, WHITE, False)],
    [("very few passes before a defensive action.", 12, WHITE, False)],
    [("", 8, WHITE, False)],
    [("Full consequence sentence.", 12, ACCENT, False)],
    [("No ambiguity about direction or meaning.", 12, ACCENT, False)],
]
rich(s, Inches(7.0), Inches(1.65), Inches(5.5), Inches(3.2), p, sp=Pt(4))

# Bottom: z-score table + key differences
rect(s, M, Inches(5.3), Inches(12), Inches(1.8), CARD)
z_levels = [
    ("\u2265 1.5  Outstanding", "1.0\u20131.5  Excellent", "0.5\u20131.0  Good",
     "\u20130.5\u20130.5  Average", "\u20131.0\u2013\u20130.5  Below avg", "< \u20131.0  Poor"),
]
p = [
    [("Z-Score \u2192 Level (same as Physical)    ", 11, MUTED, True),
     ("but we go further:", 11, ACCENT, True)],
    [("", 3, WHITE, False)],
    [("\u2022  8 metrics \u00d7 6 levels = ", 11, BODY, False),
     ("48 pre-written consequence sentences", 11, WHITE, True),
     (" \u2014 the LLM never interprets z-scores", 11, BODY, False)],
    [("\u2022  Asymmetric thresholds: strengths require z > 1.0, concerns surface at z < \u20130.5", 11, BODY, False)],
    [("\u2022  5 of 8 metrics are inverted (lower raw = better) \u2014 direction is pre-resolved", 11, BODY, False)],
    [("\u2022  Metrics grouped into ", 11, BODY, False),
     ("Strengths / Neutral / Weaknesses", 11, ACCENT, True),
     (" sections before the LLM sees them", 11, BODY, False)],
]
rich(s, Inches(0.8), Inches(5.4), Inches(11.5), Inches(1.5), p, sp=Pt(3))


# ════════════════════════════════════════════════════════════════
# SLIDES 8-11 — STEP 4: FEW-SHOT EXAMPLES
# ════════════════════════════════════════════════════════════════
SYNTH_EXAMPLES = {}
for tname in ["Liverpool", "Arsenal", "Nottingham", "Wolverhampton"]:
    profile, bp, br = team_profile(tname)
    # Build synth text like production code
    lines_str, lines_neu, lines_wk = [], [], []
    for m, lab, short, val, rank, lev, cons, z in profile:
        line = f"- {lab} \u2014 {val} \u2014 ranked {ordinal(rank)} of {N}. {cons}"
        if z > 1.0: lines_str.append(line)
        elif z < -0.5: lines_wk.append(line)
        else: lines_neu.append(line)
    synth = f"Here is a pressing profile of {tname}...\n\n"
    if lines_str: synth += "Strengths:\n" + "\n".join(lines_str) + "\n\n"
    if lines_neu: synth += "Neutral:\n" + "\n".join(lines_neu) + "\n\n"
    if lines_wk: synth += "Weaknesses:\n" + "\n".join(lines_wk) + "\n\n"
    if bp is not None:
        synth += f"Context: {tname} spent {bp:.1f}% ... ranking {ordinal(br)} out of {N}."
    SYNTH_EXAMPLES[tname] = synth

HANDWRITTEN = {
    "Liverpool": (
        "Liverpool defend with a high line and step in early when they engage, breaking opponent "
        "sequences before build-up can develop and making it very difficult to play through the "
        "defensive lines. The press wins the ball back high at a solid rate, keeps opponents locked "
        "in effectively, and makes passing under pressure noticeably harder. The clear vulnerability "
        "is behind the press: when opponents do escape, they score more often than against most "
        "teams, exposing a rest-defence that struggles to contain the damage. Overall, Liverpool "
        "trade defensive safety for front-foot disruption, backing the quality of the initial "
        "engagement to compensate for the risk when it fails."
    ),
    "Arsenal": (
        "Arsenal hold one of the highest defensive lines in the league yet press selectively rather "
        "than sustaining constant pressure, choosing specific moments to step forward. When they do "
        "engage, the coordination is strong: passing accuracy drops noticeably under their pressure "
        "and opponents find it hard to play through the lines. The press lacks volume and rarely "
        "wins the ball back in the opposition half, with few initiated sequences and limited "
        "high regains. Overall, Arsenal trade proactive disruption for a conservative, controlled "
        "engagement \u2014 high-quality pressure when it arrives but not enough of it to dominate "
        "build-up across matches."
    ),
    "Nottingham": (
        "Nottingham sit deeper than any other side in the league and press at a standard rate when "
        "they do step forward, but without the structure to back it up. Opponents escape the initial "
        "press more often than against most sides and retain their passing quality comfortably. "
        "The one genuine strength is rest-defence: when the press is bypassed, opponents rarely "
        "convert, suggesting a team that absorbs danger deeper rather than preventing it up front. "
        "Overall, Nottingham rely entirely on low-block resilience rather than pressing as a "
        "weapon, accepting territorial concession in exchange for damage limitation."
    ),
    "Wolverhampton": (
        "Wolverhampton\u2019s pressing sits close to the league norm across all dimensions \u2014 they "
        "engage at a standard frequency, step in at a typical tempo, and position themselves at a "
        "standard block height. The press does not generate standout disruption in any single area "
        "but equally does not leave the team exposed when bypassed. Opponents pass, break lines, and "
        "score after escaping at rates that sit squarely in the middle of the table. Overall, this "
        "is a neutral pressing profile \u2014 functional but unremarkable, neither a weapon nor a "
        "liability."
    ),
}

def build_fewshot_slide(tname, num):
    s = prs.slides.add_slide(blank); setup(s, num)
    tb(s, M, Inches(0.25), CW, Inches(0.5),
       f"Step 4 \u2014 Example {num-7}: {tname}", sz=28, c=WHITE, b=True)

    # Left: Synthesized Data
    rect(s, M, Inches(1.0), Inches(5.8), Inches(5.7), CARD)
    hline(s, M, Inches(1.0), Inches(5.8), MUTED, Pt(2))
    tb(s, Inches(0.8), Inches(1.1), Inches(5), Inches(0.3),
       "Synthesized Data (Input to LLM)", sz=12, c=MUTED, b=True)

    synth = SYNTH_EXAMPLES.get(tname, "")
    if len(synth) > 600: synth = synth[:600] + "..."
    tb(s, Inches(0.8), Inches(1.5), Inches(5.4), Inches(5.0),
       synth, sz=9, c=BODY, f=MONO)

    # Right: Hand-Written Output
    rect(s, Inches(6.8), Inches(1.0), Inches(5.9), Inches(5.7), CARD)
    hline(s, Inches(6.8), Inches(1.0), Inches(5.9), ACCENT, Pt(2))
    tb(s, Inches(7.0), Inches(1.1), Inches(5), Inches(0.3),
       "Our Hand-Written Output", sz=12, c=ACCENT, b=True)
    tb(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.0),
       HANDWRITTEN.get(tname, ""), sz=12, c=BODY)


build_fewshot_slide("Liverpool", 8)
build_fewshot_slide("Arsenal", 9)
build_fewshot_slide("Nottingham", 10)
build_fewshot_slide("Wolverhampton", 11)


# ════════════════════════════════════════════════════════════════
# SLIDE 12 — THE PIPELINE + KEY DECISIONS
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); setup(s, 12)
tb(s, M, Inches(0.25), CW, Inches(0.5),
   "The Pipeline", sz=30, c=WHITE, b=True)

# Pipeline boxes
boxes = [
    ("Twelve\nPressing Data", "8 raw metrics:\nchains, regains,\nPPDA, break rate..."),
    ("Statistical\nModel", "Z-scores computed\nwithin the league\n\u2192 ranks per metric"),
    ("Consequence\nStatements", "Z-score \u2192 pre-written\nfootball sentence\n(not just a label)"),
    ("Wordalization", "LLM transforms\nstructured profile into\ncoaching briefing"),
]
x = Inches(0.4)
bw = Inches(2.6)
for i, (title, desc) in enumerate(boxes):
    rect(s, x, Inches(1.1), bw, Inches(1.8), CARD)
    hline(s, x, Inches(1.1), bw, ACCENT, Pt(2))
    tb(s, x + Inches(0.15), Inches(1.2), bw - Inches(0.3), Inches(0.6),
       title, sz=13, c=WHITE, b=True, al=PP_ALIGN.CENTER)
    tb(s, x + Inches(0.15), Inches(1.8), bw - Inches(0.3), Inches(0.9),
       desc, sz=10, c=MUTED, al=PP_ALIGN.CENTER)
    if i < len(boxes) - 1:
        tb(s, x + bw, Inches(1.7), Inches(0.5), Inches(0.4),
           "\u2192", sz=20, c=DIM, al=PP_ALIGN.CENTER)
    x += bw + Inches(0.55)

# Key Decisions
rect(s, M, Inches(3.3), Inches(12), Inches(3.7), CARD)
tb(s, Inches(0.8), Inches(3.4), Inches(5), Inches(0.3),
   "Key Decisions", sz=16, c=WHITE, b=True)
hline(s, Inches(0.8), Inches(3.7), Inches(3), ACCENT, Pt(2))

decisions = [
    ("\u2022  Consequence sentences, not word labels",
     "Physical says \u2018poor in Speed\u2019 \u2014 we say \u2018opponents pass comfortably under this pressure\u2019"),
    ("\u2022  Pre-grouped Strengths / Neutral / Weaknesses",
     "The LLM doesn\u2019t decide what\u2019s good or bad \u2014 we do"),
    ("\u2022  Block % as separate context layer",
     "Defensive positioning is injected without distorting the metric set"),
    ("\u2022  4 diverse few-shot examples",
     "Aggressive (Liverpool), selective (Arsenal), fragile (Nottingham), neutral (Wolverhampton)"),
    ("\u2022  Grounded chat with RAG",
     "Follow-up answers use synthesized text + embedding search (top 5, cosine > 0.7)"),
]
y = Inches(3.85)
for title, desc in decisions:
    tb(s, Inches(0.8), y, Inches(11.5), Inches(0.25), title, sz=12, c=WHITE, b=True)
    tb(s, Inches(1.2), y + Inches(0.25), Inches(11), Inches(0.25), desc, sz=10, c=MUTED)
    y += Inches(0.58)


# ════════════════════════════════════════════════════════════════
# SLIDE 13 — BLOCK % CONTEXT (BONUS)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); setup(s, 13)
tb(s, M, Inches(0.25), CW, Inches(0.5),
   "Bonus: Block Height Context", sz=30, c=WHITE, b=True)
tb(s, M, Inches(0.8), CW, Inches(0.3),
   "Something Physical Analyst doesn\u2019t have \u2014 tactical context beyond the metrics", sz=14, c=MUTED)

all_b = df_league["Block %"]
block_sorted = df_league[["Team", "Block %"]].sort_values("Block %", ascending=False)

# Explanation
rect(s, M, Inches(1.4), Inches(7), Inches(5.2), CARD)
hline(s, M, Inches(1.4), Inches(7), ACCENT, Pt(2))
p = [
    [("Why add this?", 16, WHITE, True)],
    [("", 4, WHITE, False)],
    [("Block % tells the LLM ", 13, BODY, False), ("where", 13, WHITE, True),
     (" the team defends \u2014 not ", 13, BODY, False), ("how well.", 13, WHITE, True)],
    [("A deep team pressing poorly means something different from a high-line team pressing poorly.", 12, MUTED, False)],
    [("", 8, WHITE, False)],
    [("How we handle it:", 14, WHITE, True)],
    [("", 3, WHITE, False)],
    [("\u2022  Loaded from a separate CSV \u2014 not mixed with pressing metrics", 12, BODY, False)],
    [("\u2022  Not z-scored or ranked as a quality metric", 12, BODY, False)],
    [("\u2022  Explicitly framed: \u201cThis is not a quality judgement\u201d", 12, BODY, False)],
    [("\u2022  Few-shot examples show the LLM how to weave it in naturally", 12, BODY, False)],
    [("", 8, WHITE, False)],
    [("League:  ", 12, MUTED, False),
     (f"Min {all_b.min():.1f}%", 12, RED, True),
     (f"   Median {all_b.median():.1f}%", 12, BODY, True),
     (f"   Max {all_b.max():.1f}%", 12, ACCENT, True)],
]
rich(s, Inches(0.8), Inches(1.55), Inches(6.6), Inches(4.8), p, sp=Pt(3))

# Right: top/bottom
rect(s, Inches(7.8), Inches(1.4), Inches(4.9), Inches(2.3), CARD)
hline(s, Inches(7.8), Inches(1.4), Inches(4.9), ACCENT, Pt(2))
top_p = [[("Highest Block %", 13, ACCENT, True)], [("", 3, WHITE, False)]]
for _, row in block_sorted.head(5).iterrows():
    top_p.append([(f"  {row['Team']:<20} {row['Block %']:.1f}%", 11, BODY, False, MONO)])
rich(s, Inches(8.0), Inches(1.55), Inches(4.5), Inches(2.0), top_p, sp=Pt(3))

rect(s, Inches(7.8), Inches(3.95), Inches(4.9), Inches(2.3), CARD)
hline(s, Inches(7.8), Inches(3.95), Inches(4.9), RED, Pt(2))
bot_p = [[("Lowest Block %", 13, RED, True)], [("", 3, WHITE, False)]]
for _, row in block_sorted.tail(5).iterrows():
    bot_p.append([(f"  {row['Team']:<20} {row['Block %']:.1f}%", 11, BODY, False, MONO)])
rich(s, Inches(8.0), Inches(4.1), Inches(4.5), Inches(2.0), bot_p, sp=Pt(3))


# ── Save ────────────────────────────────────────────────────────
OUT = "Pressing_Analyst_slides_vf.pptx"
prs.save(OUT)
print(f"Done! {OUT} \u2014 {len(prs.slides)} slides")
