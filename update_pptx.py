"""
Update Pressing_Analyst_final.pptx based on code-review feedback.
Produces Pressing_Analyst_v2.pptx.
"""

from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

SRC = r"d:\Twelve\Twelve-GPT-Educational\Pressing_Analyst_final.pptx"
DST = r"d:\Twelve\Twelve-GPT-Educational\Pressing_Analyst_v2.pptx"

ORANGE = RGBColor(0xFF, 0x64, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0xA0, 0xA0, 0xA0)
LGRAY  = RGBColor(0xD2, 0xD2, 0xD2)
GREEN  = RGBColor(0x00, 0xC8, 0x78)
AMBER  = RGBColor(0xFF, 0xB7, 0x00)

LEFT_MARGIN = Emu(548640)
BODY_LEFT   = Emu(914400)
BODY_WIDTH  = Emu(10058400)
TITLE_WIDTH = Emu(10972800)

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _style_run(run, sz, bold, color):
    rPr = run._r.find(f"{{{A_NS}}}rPr")
    if rPr is None:
        rPr = etree.SubElement(run._r, f"{{{A_NS}}}rPr")
        run._r.insert(0, rPr)
    rPr.set("sz", str(sz))
    rPr.set("lang", "en-GB")
    rPr.set("dirty", "0")
    if bold:
        rPr.set("b", "1")
    fill = rPr.find(f"{{{A_NS}}}solidFill")
    if fill is None:
        fill = etree.SubElement(rPr, f"{{{A_NS}}}solidFill")
    for ch in list(fill):
        fill.remove(ch)
    clr_el = etree.SubElement(fill, f"{{{A_NS}}}srgbClr")
    clr_el.set("val", str(color))


def _style_paragraph_default(paragraph, sz, bold, color):
    pPr = paragraph._p.find(f"{{{A_NS}}}pPr")
    if pPr is None:
        pPr = etree.SubElement(paragraph._p, f"{{{A_NS}}}pPr")
        paragraph._p.insert(0, pPr)
    pPr.set("algn", "l")
    defrp = pPr.find(f"{{{A_NS}}}defRPr")
    if defrp is None:
        defrp = etree.SubElement(pPr, f"{{{A_NS}}}defRPr")
    defrp.set("sz", str(sz))
    defrp.set("b", "1" if bold else "0")
    fill = defrp.find(f"{{{A_NS}}}solidFill")
    if fill is None:
        fill = etree.SubElement(defrp, f"{{{A_NS}}}solidFill")
    for ch in list(fill):
        fill.remove(ch)
    clr_el = etree.SubElement(fill, f"{{{A_NS}}}srgbClr")
    clr_el.set("val", str(color))


def make_textbox(slide, left, top, width, height, text, sz=1300, bold=False, color=WHITE):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _style_paragraph_default(p, sz, bold, color)
    _style_run(run, sz, bold, color)
    return txbox


def add_multiline_textbox(slide, left, top, width, height, lines):
    """lines = [(text, sz, bold, color), ...]"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    bodyPr = tf._txBody.find(f"{{{A_NS}}}bodyPr")
    for ch in list(bodyPr):
        bodyPr.remove(ch)
    etree.SubElement(bodyPr, f"{{{A_NS}}}spAutoFit")

    for i, (text, sz, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            etree.SubElement(tf._txBody, f"{{{A_NS}}}p")
            p = tf.paragraphs[-1]
        run = p.add_run()
        run.text = text
        _style_paragraph_default(p, sz, bold, color)
        _style_run(run, sz, bold, color)
    return txbox


def create_blank_slide(prs):
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            return prs.slides.add_slide(layout)
    return prs.slides.add_slide(prs.slide_layouts[6])


def reorder_slides(prs, new_order):
    """new_order = list of 0-based current indices in desired order."""
    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    ordered = [items[i] for i in new_order]
    for item in list(sldIdLst):
        sldIdLst.remove(item)
    for item in ordered:
        sldIdLst.append(item)


# ── Load ──────────────────────────────────────────
prs = Presentation(SRC)

# Current 16 slides (0-indexed):
# 0:  Title
# 1:  What to measure
# 2:  Wall or gamble
# 3:  Data To Text
# 4:  System Prompt
# 5:  Q&A Pairs
# 6:  Synthesize (Liverpool)
# 7:  Few Shot Example
# 8:  Prompt Instruction
# 9:  Pipeline
# 10: Fulham
# 11: Bournemouth
# 12: Bournemouth DUPLICATE
# 13: Ipswich
# 14: Arsenal
# 15: Q&A closing

# ═══════════════════════════════════════════════════
# 1. Slide 5 (idx 4) — System Prompt: add alignment trick
# ═══════════════════════════════════════════════════
s5 = prs.slides[4]
make_textbox(s5, BODY_LEFT, Emu(3200000), BODY_WIDTH, Emu(274320),
             "Terminology Alignment", sz=1100, bold=True, color=GREEN)
make_textbox(s5, BODY_LEFT, Emu(3520000), BODY_WIDTH, Emu(274320),
             'User: "Do you refer to the game as soccer or football?"',
             sz=1100, bold=False, color=LGRAY)
make_textbox(s5, BODY_LEFT, Emu(3800000), BODY_WIDTH, Emu(274320),
             'Assistant: "I refer to the game as football. I write for coaches, not commentators."',
             sz=1100, bold=False, color=LGRAY)
make_textbox(s5, BODY_LEFT, Emu(4150000), BODY_WIDTH, Emu(400000),
             "This single exchange locks in British English terminology and coaching register "
             "before any data is processed.",
             sz=1000, bold=False, color=GRAY)

# ═══════════════════════════════════════════════════
# 2. Slide 6 (idx 5) — Q&A Pairs: add dual-purpose note
# ═══════════════════════════════════════════════════
s6 = prs.slides[5]
add_multiline_textbox(s6, LEFT_MARGIN, Emu(6200000), TITLE_WIDTH, Emu(500000), [
    ("Dual purpose: ", 1000, True, GREEN),
    ("These Q&A pairs are used for few-shot learning during description generation "
     "AND for semantic retrieval (RAG) during the interactive chat. "
     "When a user asks a follow-up, the 5 most relevant pairs are retrieved "
     "via cosine similarity and injected into the prompt.",
     1000, False, GRAY),
])

# ═══════════════════════════════════════════════════
# 3. Slide 7 (idx 6) — Synthesize: context-block note
# ═══════════════════════════════════════════════════
s7 = prs.slides[6]
add_multiline_textbox(s7, LEFT_MARGIN, Emu(5700000), TITLE_WIDTH, Emu(600000), [
    ("Why is Context separate? ", 1000, True, AMBER),
    ("High/medium block % is a tactical preference, not a quality metric. "
     "It is shown outside Strengths/Weaknesses so the LLM treats it as background, "
     "not as something good or bad.",
     1000, False, GRAY),
])

# ═══════════════════════════════════════════════════
# 4. Slide 8 (idx 7) — Few Shot: 4-profile strategy
# ═══════════════════════════════════════════════════
s8 = prs.slides[7]
add_multiline_textbox(s8, LEFT_MARGIN, Emu(5500000), TITLE_WIDTH, Emu(1100000), [
    ("Few-shot strategy: 4 contrasting profiles", 1000, True, GREEN),
    ("Liverpool \u2014 aggressive press, weak rest-defence   |   "
     "Arsenal \u2014 selective, low volume   |   "
     "Nottingham \u2014 deep block, fragile press, strong deep defence   |   "
     "Wolverhampton \u2014 entirely neutral baseline",
     900, False, LGRAY),
    ("Each archetype teaches the model a different pressing identity. "
     "Without this diversity the model defaults to one lens.",
     900, False, GRAY),
])

# ═══════════════════════════════════════════════════
# 5. NEW: Guardrails & Constraints slide (appended, will reorder later)
# ═══════════════════════════════════════════════════
gs = create_blank_slide(prs)  # idx 16
make_textbox(gs, LEFT_MARGIN, Emu(365760), TITLE_WIDTH, Emu(548640),
             "Guardrails & Constraints", sz=2800, bold=True, color=ORANGE)
make_textbox(gs, LEFT_MARGIN, Emu(1000000), TITLE_WIDTH, Emu(365760),
             "Design decisions that keep the output reliable",
             sz=1400, bold=False, color=GRAY)
add_multiline_textbox(gs, BODY_LEFT, Emu(1500000), BODY_WIDTH, Emu(5200000), [
    ("\u2776  Asymmetric Grouping", 1300, True, GREEN),
    ("Strengths require z > 1.0. Weaknesses surface at z < \u22120.5 \u2014 deliberately earlier. "
     "Coaches need to catch problems before they become crises.", 1100, False, LGRAY),
    ("", 500, False, LGRAY),
    ("\u2777  Anti-Hallucination", 1300, True, GREEN),
    ("\u201CDo not invent consequences not in the data "
     "(e.g. rapid restarts, transition speed, possession recycling).\u201D", 1100, False, LGRAY),
    ("", 500, False, LGRAY),
    ("\u2778  No Jargon Leakage", 1300, True, GREEN),
    ("\u201CDo not include metric names, level labels (e.g. excellent, poor), "
     "or parenthetical references in your output.\u201D", 1100, False, LGRAY),
    ("", 500, False, LGRAY),
    ("\u2779  Structured Output", 1300, True, GREEN),
    ("Exactly 4 sentences: pressing identity \u2192 strengths \u2192 weaknesses \u2192 trade-off.",
     1100, False, LGRAY),
    ("", 500, False, LGRAY),
    ("\u277A  Token Limit", 1300, True, GREEN),
    ("250 tokens max. A coach reads it in under 30 seconds.", 1100, False, LGRAY),
    ("", 500, False, LGRAY),
    ("\u277B  Inverse Metrics Resolved Upstream", 1300, True, GREEN),
    ("5 of 8 metrics are inverted (lower = better). Direction is resolved via "
     "pre-written consequence sentences \u2014 the LLM never interprets polarity.",
     1100, False, LGRAY),
])

# ═══════════════════════════════════════════════════
# 6. NEW: Chat Pipeline slide (appended, will reorder later)
# ═══════════════════════════════════════════════════
cs = create_blank_slide(prs)  # idx 17
make_textbox(cs, LEFT_MARGIN, Emu(365760), TITLE_WIDTH, Emu(548640),
             "Interactive Chat Pipeline", sz=2800, bold=True, color=ORANGE)
make_textbox(cs, LEFT_MARGIN, Emu(1000000), TITLE_WIDTH, Emu(365760),
             "How the system answers follow-up questions after the initial briefing",
             sz=1400, bold=False, color=GRAY)
add_multiline_textbox(cs, BODY_LEFT, Emu(1500000), BODY_WIDTH, Emu(5200000), [
    ("\u2776  User Question", 1300, True, GREEN),
    ("User types a follow-up in the chat interface (max 500 characters).",
     1100, False, LGRAY),
    ("", 500, False, LGRAY),
    ("\u2777  Embed & Search", 1300, True, GREEN),
    ("Question is embedded and compared against the Q&A library via cosine similarity. "
     "Top 5 most relevant pairs are retrieved.", 1100, False, LGRAY),
    ("", 500, False, LGRAY),
    ("\u2778  Build Context", 1300, True, GREEN),
    ("Retrieved Q&A pairs + cached team pressing profile (synthesized text) "
     "are combined into one context block.", 1100, False, LGRAY),
    ("", 500, False, LGRAY),
    ("\u2779  Chat System Prompt", 1300, True, GREEN),
    ("\u201CYou are a tactical data analyst... provide 2-sentence answers. "
     "Do not deviate from the information provided.\u201D", 1100, False, LGRAY),
    ("", 500, False, LGRAY),
    ("\u277A  LLM Call (temperature = 0.3)", 1300, True, GREEN),
    ("Low temperature for factual, consistent output. "
     "Same question about the same team = same answer every time.", 1100, False, LGRAY),
    ("", 500, False, LGRAY),
    ("\u277B  Design Note", 1300, True, AMBER),
    ("The same Q&A pairs serve both few-shot learning (description) and "
     "RAG retrieval (chat). One knowledge base, two functions.", 1100, False, GRAY),
])

# ═══════════════════════════════════════════════════
# 7. Reorder all slides to final sequence
# ═══════════════════════════════════════════════════
# Current indices after appending two new slides:
#  0: Title
#  1: What to measure
#  2: Wall or gamble
#  3: Data To Text
#  4: System Prompt (updated)
#  5: Q&A Pairs (updated)
#  6: Synthesize Liverpool (updated)
#  7: Few Shot Example (updated)
#  8: Prompt Instruction
#  9: Pipeline
# 10: Fulham
# 11: Bournemouth
# 12: Bournemouth DUPLICATE  <-- skip this
# 13: Ipswich
# 14: Arsenal
# 15: Q&A closing
# 16: Guardrails (NEW)
# 17: Chat Pipeline (NEW)

desired_order = [
    0,   # 1. Title
    1,   # 2. What to measure
    2,   # 3. Wall or gamble
    3,   # 4. Data To Text
    4,   # 5. System Prompt (+ alignment trick)
    5,   # 6. Q&A Pairs (+ dual purpose note)
    6,   # 7. Synthesize Liverpool (+ context note)
    7,   # 8. Few Shot Example (+ 4 profiles)
    8,   # 9. Prompt Instruction
    9,   # 10. Pipeline
    16,  # 11. Guardrails & Constraints (NEW)
    10,  # 12. Fulham
    11,  # 13. Bournemouth
    13,  # 14. Ipswich (skip duplicate Bournemouth at 12)
    14,  # 15. Arsenal
    17,  # 16. Chat Pipeline (NEW)
    15,  # 17. Q&A closing
]

reorder_slides(prs, desired_order)

# ═══════════════════════════════════════════════════
prs.save(DST)
print(f"Saved to {DST}")
print(f"Total slides: {len(list(prs.slides._sldIdLst))}")
