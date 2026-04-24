"""
Build Pressing Analyst Storyboard — styled after the Corner Analyst storyboard deck.
Five slides: 1 framing slide + 4 question slides.

Run:  python build_pressing_storyboard.py
Out:  Pressing_Analyst_Storyboard.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette (light theme, Corner Analyst style) ──────────────────
BG        = RGBColor(255, 255, 255)
CARD_BG   = RGBColor(246, 247, 249)
NAVY      = RGBColor(26, 46, 74)
NAVY_DARK = RGBColor(15, 30, 50)
GREEN     = RGBColor(34, 181, 115)
BLUE      = RGBColor(26, 77, 143)
RED       = RGBColor(214, 69, 65)
AMBER     = RGBColor(245, 158, 11)
TEXT      = RGBColor(30, 35, 45)
MUTED     = RGBColor(102, 112, 128)
DIM       = RGBColor(160, 170, 180)
LINE      = RGBColor(226, 229, 233)
WHITE     = RGBColor(255, 255, 255)
PITCH_BG  = RGBColor(28, 45, 70)
FONT      = "Calibri"
MONO      = "Consolas"

# ── Dimensions (16:9) ────────────────────────────────────────────
SW = Inches(13.333)
SH = Inches(7.5)
M  = Inches(0.55)
CW = SW - 2 * M

TOTAL_SLIDES = 5


# ── Shape helpers ────────────────────────────────────────────────
def _noborder(shape):
    shape.line.fill.background()


def rect(slide, l, t, w, h, fill, rounded=True, radius=0.04):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    _noborder(s)
    if rounded and s.adjustments:
        s.adjustments[0] = radius
    return s


def hline(slide, l, t, w, color=LINE, th=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, th)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    _noborder(s)
    return s


def vline(slide, l, t, h, color, th=Pt(3)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, th, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    _noborder(s)
    return s


def tb(slide, l, t, w, h, text, sz=13, c=TEXT, b=False, al=PP_ALIGN.LEFT, f=FONT,
       anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.name = f
    p.font.bold = b
    p.alignment = al
    return box


def rich(slide, l, t, w, h, paras, al=PP_ALIGN.LEFT, sp=Pt(4), anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for pi, runs in enumerate(paras):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = al
        p.space_after = sp
        for text, sz, clr, bld, *extra in runs:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(sz)
            r.font.color.rgb = clr
            r.font.name = extra[0] if extra else FONT
            r.font.bold = bld
    return box


# ── Slide chrome ─────────────────────────────────────────────────
def chrome(slide, section_tag, page_num, footer="Pressing Analyst Storyboard"):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    tb(slide, M, Inches(0.35), Inches(6), Inches(0.28),
       section_tag.upper(), sz=9, c=MUTED, b=True)
    tb(slide, SW - Inches(2.0), Inches(0.35), Inches(1.5), Inches(0.28),
       "TWELVE", sz=14, c=NAVY, b=True, al=PP_ALIGN.RIGHT)
    hline(slide, M, SH - Inches(0.45), CW, LINE, Pt(0.5))
    tb(slide, M, SH - Inches(0.38), Inches(6), Inches(0.25),
       footer, sz=9, c=MUTED)
    tb(slide, SW - Inches(1.5), SH - Inches(0.38), Inches(1), Inches(0.25),
       f"{page_num:02d}", sz=9, c=MUTED, b=True, al=PP_ALIGN.RIGHT)


def top_card(slide, l, t, w, h, title, title_color, body_lines):
    """Card with a thin colored top border — used in the framing slide."""
    rect(slide, l, t, w, h, CARD_BG, radius=0.03)
    hline(slide, l, t, w, title_color, Pt(3))
    tb(slide, l + Inches(0.2), t + Inches(0.22), w - Inches(0.4), Inches(0.35),
       title, sz=16, c=NAVY, b=True)
    paras = [[(line, 10, MUTED, False)] for line in body_lines]
    rich(slide, l + Inches(0.2), t + Inches(0.7), w - Inches(0.4), h - Inches(0.8),
         paras, sp=Pt(3))


def left_card(slide, l, t, w, h, title, title_color, body):
    """Card with a thick colored left strip — used in question slides."""
    rect(slide, l, t, w, h, CARD_BG, radius=0.04)
    vline(slide, l, t, h, title_color, Pt(4))
    tb(slide, l + Inches(0.2), t + Inches(0.15), w - Inches(0.3), Inches(0.3),
       title, sz=13, c=NAVY, b=True)
    tb(slide, l + Inches(0.2), t + Inches(0.5), w - Inches(0.3), h - Inches(0.6),
       body, sz=10, c=MUTED)


def stat_card(slide, l, t, w, h, label, value, value_color=RED):
    rect(slide, l, t, w, h, CARD_BG, radius=0.05)
    tb(slide, l, t + Inches(0.12), w, Inches(0.3),
       label, sz=10, c=MUTED, b=True, al=PP_ALIGN.CENTER)
    tb(slide, l, t + Inches(0.45), w, Inches(0.5),
       value, sz=24, c=value_color, b=True, al=PP_ALIGN.CENTER)


def evidence_card(slide, l, t, w, h, title, bullets, accent=GREEN):
    rect(slide, l, t, w, h, CARD_BG, radius=0.04)
    vline(slide, l, t, h, accent, Pt(3))
    tb(slide, l + Inches(0.2), t + Inches(0.12), w - Inches(0.3), Inches(0.3),
       title, sz=11, c=NAVY, b=True)
    paras = [[(f"{i+1}. {txt}" if False else txt, 9, MUTED, False)]
             for i, txt in enumerate(bullets)]
    rich(slide, l + Inches(0.2), t + Inches(0.5), w - Inches(0.3), h - Inches(0.55),
         paras, sp=Pt(3))


def _emu(x):
    """Coerce any numeric (int / float / Emu) to an Emu-typed int for pptx."""
    return Emu(int(x))


def supporting_row(slide, l, t, w, label, fill_pct, accent=GREEN):
    """Small row: label + progress bar, used in supporting visuals card."""
    tb(slide, l, t, Inches(1.8), Inches(0.2), label, sz=9, c=MUTED)
    bar_l = l + Inches(1.9)
    bar_w = _emu(w - Inches(1.9))
    rect(slide, bar_l, t + Inches(0.04), bar_w, Inches(0.12), LINE, radius=0.5)
    if fill_pct > 0:
        rect(slide, bar_l, t + Inches(0.04), _emu(bar_w * fill_pct),
             Inches(0.12), accent, radius=0.5)


def pitch_mock(slide, l, t, w, h, caption=""):
    """Dark navy pitch-like placeholder for the main visual slot."""
    rect(slide, l, t, w, h, PITCH_BG, radius=0.02)
    inner_l = l + Inches(0.3)
    inner_t = t + Inches(0.3)
    inner_w = _emu(w - Inches(0.6))
    inner_h = _emu(h - Inches(0.6))
    for i in range(1, 4):
        x = _emu(inner_l + inner_w * i / 4)
        rect(slide, x, inner_t, Pt(0.5), inner_h,
             RGBColor(60, 80, 110), rounded=False)
    mid_x = _emu(inner_l + inner_w / 2 - Inches(0.02))
    rect(slide, mid_x, inner_t, Inches(0.04), inner_h, WHITE, rounded=False)
    if caption:
        tb(slide, l, _emu(t + h - Inches(0.35)), w, Inches(0.25),
           caption, sz=9, c=RGBColor(180, 190, 210), al=PP_ALIGN.CENTER)


def dot_markers(slide, l, t, w, h, positions, color=GREEN, size=Inches(0.22)):
    """Draw small dots on top of the pitch mock, positions are (x_frac, y_frac)."""
    for xf, yf in positions:
        cx = _emu(l + w * xf - size / 2)
        cy = _emu(t + h * yf - size / 2)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, size, size)
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        _noborder(dot)


# ── Build ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
blank = prs.slide_layouts[6]


# ═════════════════════════════════════════════════════════════════
# SLIDE 01 — FRAMING: What the Pressing Analyst should do
# ═════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
chrome(s, "Analyst Framing", 1)

tb(s, M, Inches(0.85), CW, Inches(0.7),
   "What the Pressing Analyst should do", sz=32, c=NAVY, b=True)
tb(s, M, Inches(1.55), CW - Inches(0.5), Inches(0.8),
   "The storyboard should position the assistant as a question-led tactical analyst, "
   "not just a text generator. Each page should connect a pressing question to clear "
   "evidence already present in the current page.",
   sz=11, c=MUTED)

card_t = Inches(2.55)
card_h = Inches(1.65)
card_w = Inches(4.05)
gap    = Inches(0.15)

top_card(s, M, card_t, card_w, card_h,
         "Analyst promise", GREEN,
         [
             "Translate pressing data into",
             "coach-friendly tactical answers. Stay",
             "concise, compare against league",
             "context and explain intensity curve,",
             "bypass lanes and genuine cost.",
         ])

top_card(s, M + card_w + gap, card_t, card_w, card_h,
         "Evidence already available", BLUE,
         [
             "PPDA, pressing chains, regains in",
             "opp half, pass accuracy under",
             "pressure, press-break rate, line-",
             "breaking pass rate, goals after",
             "broken press, block %.",
         ])

top_card(s, M + 2 * (card_w + gap), card_t, card_w, card_h,
         "Output style already aligned", RED,
         [
             "The current prompt already produces",
             "analyst-grade copy: no raw metric",
             "labels, short tactical summaries,",
             "profile \u2192 reward \u2192 risk \u2192 closing",
             "cause-and-effect language.",
         ])

tb(s, M, Inches(4.5), Inches(6), Inches(0.4),
   "Recommended questions", sz=16, c=NAVY, b=True)

questions = [
    "1. How does the team\u2019s pressing intensity change within matches and across the season?",
    "2. Which opposition-half zones generate the most ball recoveries?",
    "3. Which lanes does the opponent use to break the press?",
    "4. Which zones do opponents find shots from after breaking the press?",
]
y = Inches(4.95)
for q in questions:
    tb(s, M + Inches(0.1), y, Inches(7.2), Inches(0.3),
       q, sz=11, c=TEXT)
    y += Inches(0.32)

sp_l = M + Inches(7.5)
sp_t = Inches(4.6)
sp_w = CW - Inches(7.5)
sp_h = Inches(1.9)
rect(s, sp_l, sp_t, sp_w, sp_h, NAVY_DARK, radius=0.03)
vline(s, sp_l, sp_t, sp_h, GREEN, Pt(4))
tb(s, sp_l + Inches(0.25), sp_t + Inches(0.2), sp_w - Inches(0.4), Inches(0.35),
   "Storyboard principle", sz=14, c=WHITE, b=True)
tb(s, sp_l + Inches(0.25), sp_t + Inches(0.65), sp_w - Inches(0.4), sp_h - Inches(0.8),
   "Each page should show the question, the answer shape, the core visual "
   "evidence and the tactical decision it unlocks.",
   sz=11, c=RGBColor(200, 210, 225))


# ═════════════════════════════════════════════════════════════════
# Helper for question slides
# ═════════════════════════════════════════════════════════════════
def build_question_slide(num, q_tag, title, subtitle,
                          user_asks, analyst_answer, decision_value,
                          stat_cards, main_caption, pitch_dots,
                          supporting_title, supporting_rows,
                          key_evidence, also_shown):
    s = prs.slides.add_slide(blank)
    chrome(s, q_tag, num)

    tb(s, M, Inches(0.85), CW, Inches(0.7),
       title, sz=28, c=NAVY, b=True)
    tb(s, M, Inches(1.5), CW - Inches(0.5), Inches(0.6),
       subtitle, sz=11, c=MUTED)

    # Left column — 3 cards
    lc_w = Inches(3.15)
    lc_l = M
    lc_t = Inches(2.35)
    lc_h = Inches(1.3)
    lc_gap = Inches(0.2)

    left_card(s, lc_l, lc_t, lc_w, lc_h, "User asks", GREEN, user_asks)
    left_card(s, lc_l, lc_t + lc_h + lc_gap, lc_w, lc_h,
              "Analyst answer", BLUE, analyst_answer)
    left_card(s, lc_l, lc_t + 2 * (lc_h + lc_gap), lc_w, lc_h,
              "Decision value", RED, decision_value)

    # Right column — storyboard frame
    sf_l = lc_l + lc_w + Inches(0.3)
    sf_t = Inches(2.35)
    sf_w = _emu(SW - sf_l - M)
    sf_h = Inches(4.1)
    rect(s, sf_l, sf_t, sf_w, sf_h, WHITE, radius=0.02)
    hline(s, sf_l, sf_t, sf_w, LINE, Pt(0.75))
    hline(s, sf_l, _emu(sf_t + sf_h), sf_w, LINE, Pt(0.75))
    vline(s, sf_l, sf_t, sf_h, LINE, Pt(0.75))
    vline(s, _emu(sf_l + sf_w), sf_t, sf_h, LINE, Pt(0.75))

    tb(s, sf_l + Inches(0.2), sf_t + Inches(0.12), Inches(3), Inches(0.3),
       "Storyboard frame", sz=10, c=MUTED, b=True)
    tb(s, sf_l + Inches(0.2), sf_t + Inches(0.38), _emu(sf_w - Inches(0.4)), Inches(0.3),
       main_caption, sz=9, c=DIM)

    # Stat cards row (top)
    n_stats = len(stat_cards)
    stat_l = sf_l + Inches(0.25)
    stat_t = sf_t + Inches(0.75)
    stat_w = _emu((sf_w - Inches(0.5)) / n_stats)
    stat_h = Inches(0.9)
    for i, (label, value) in enumerate(stat_cards):
        stat_card(s, _emu(stat_l + i * stat_w + Inches(0.04)),
                  stat_t, _emu(stat_w - Inches(0.08)), stat_h,
                  label, value, RED)

    # Main visual (pitch) + supporting visuals side-by-side
    mv_l = stat_l
    mv_t = _emu(stat_t + stat_h + Inches(0.15))
    mv_w = _emu(sf_w * 0.62)
    mv_h = _emu(sf_h - (mv_t - sf_t) - Inches(0.25))
    pitch_mock(s, mv_l, mv_t, mv_w, mv_h, main_caption)
    if pitch_dots:
        dot_markers(s, mv_l, mv_t, mv_w, mv_h, pitch_dots, GREEN)

    sv_l = _emu(mv_l + mv_w + Inches(0.2))
    sv_t = mv_t
    sv_w = _emu(sf_w - (sv_l - sf_l) - Inches(0.25))
    sv_h = mv_h
    rect(s, sv_l, sv_t, sv_w, sv_h, CARD_BG, radius=0.04)
    tb(s, sv_l + Inches(0.15), sv_t + Inches(0.1), _emu(sv_w - Inches(0.2)), Inches(0.25),
       supporting_title, sz=10, c=NAVY, b=True)
    row_t = _emu(sv_t + Inches(0.45))
    for label, pct in supporting_rows:
        supporting_row(s, sv_l + Inches(0.15), row_t,
                       _emu(sv_w - Inches(0.3)), label, pct, GREEN)
        row_t = _emu(row_t + Inches(0.3))

    # Key evidence + Also shown bottom cards
    bc_t = _emu(sf_t + sf_h + Inches(0.2))
    bc_h = Inches(0.95)
    bc_w = _emu((sf_w - Inches(0.3)) / 2)
    evidence_card(s, sf_l, bc_t, bc_w, bc_h, "Key evidence", key_evidence, GREEN)
    evidence_card(s, _emu(sf_l + bc_w + Inches(0.3)), bc_t, bc_w, bc_h,
                  "Also shown", also_shown, BLUE)


# ═════════════════════════════════════════════════════════════════
# SLIDE 02 — Q1: Pressing intensity over time
# ═════════════════════════════════════════════════════════════════
build_question_slide(
    num=2,
    q_tag="Question 01",
    title="How does pressing intensity change within matches and across the season?",
    subtitle="This question defines pressing dynamics on two scales: within the 90 minutes and across the season.",
    user_asks="Is this a 90-minute pressing side, and does that intensity hold up across the season?",
    analyst_answer="Name the within-match curve (sustained / front-loaded / fades late) and the season trend (rising / flat / declining).",
    decision_value="Substitution windows for the match layer; form and fixture-load management for the season layer.",
    stat_cards=[
        ("First 15\u2019 PPDA", "7.2"),
        ("Last 15\u2019 PPDA", "14.8"),
        ("Season trend", "\u2197"),
        ("High-press matches", "68%"),
    ],
    main_caption="Within-match rolling PPDA curve with league-median band",
    pitch_dots=[(0.15, 0.35), (0.30, 0.40), (0.45, 0.48), (0.60, 0.58),
                (0.75, 0.65), (0.85, 0.70)],
    supporting_title="Supporting visuals",
    supporting_rows=[
        ("Chains per match", 0.80),
        ("Regains opp half", 0.65),
        ("Season PPDA trend", 0.55),
    ],
    key_evidence=[
        "Within-match 5-min rolling PPDA vs league median.",
        "Match-by-match PPDA timeline across the season.",
    ],
    also_shown=[
        "Game-state split (leading / drawing / trailing) PPDA.",
        "High-press minutes share by match phase.",
    ],
)


# ═════════════════════════════════════════════════════════════════
# SLIDE 03 — Q2: Recovery zones
# ═════════════════════════════════════════════════════════════════
build_question_slide(
    num=3,
    q_tag="Question 02",
    title="Which opposition-half zones generate the most ball recoveries?",
    subtitle="This question defines recovery identity: dominant regain channel and regain source.",
    user_asks="Where on the pitch do they actually win the ball back, and does that match their block?",
    analyst_answer="Name the dominant recovery zone and whether it comes from pressing triggers or opp build-up mistakes.",
    decision_value="Reveals which channel to attack in transition and which opposition build-up pattern is vulnerable.",
    stat_cards=[
        ("Left half", "36%"),
        ("Central", "28%"),
        ("Right half", "36%"),
        ("Under 30m from goal", "41%"),
    ],
    main_caption="Opposition-half recovery heatmap — 6-zone grid",
    pitch_dots=[(0.20, 0.30), (0.20, 0.70),
                (0.50, 0.35), (0.50, 0.65),
                (0.80, 0.32), (0.80, 0.68)],
    supporting_title="Supporting visuals",
    supporting_rows=[
        ("Tackle regains",      0.55),
        ("Interception regains", 0.70),
        ("Turnover regains",    0.35),
    ],
    key_evidence=[
        "Zone-level regains vs league average.",
        "Pressing-chain start vs recovery end mapping.",
    ],
    also_shown=[
        "Regain-to-shot conversion by zone.",
        "Dominant recovery channel marker.",
    ],
)


# ═════════════════════════════════════════════════════════════════
# SLIDE 04 — Q3: Break lanes
# ═════════════════════════════════════════════════════════════════
build_question_slide(
    num=4,
    q_tag="Question 03",
    title="Which lanes does the opponent use to break the press?",
    subtitle="This question defines vulnerability identity: the primary bypass lane and its first-pass type.",
    user_asks="Which lane do opponents use to beat our press — central, wide or by going long?",
    analyst_answer="Name the primary bypass lane and the first-pass type that opens it (short through central, wide progression, long ball).",
    decision_value="Exposes which pressing unit is being beaten, informs press-shape and cover drills.",
    stat_cards=[
        ("Central lane", "34%"),
        ("Left channel", "24%"),
        ("Right channel", "27%"),
        ("Long-ball bypass", "15%"),
    ],
    main_caption="Lane-split break-the-press map — 3 vertical corridors + long-ball overlay",
    pitch_dots=[(0.35, 0.50), (0.45, 0.50), (0.55, 0.50),
                (0.20, 0.25), (0.20, 0.75)],
    supporting_title="Supporting visuals",
    supporting_rows=[
        ("Line-break pass rate",   0.72),
        ("Pass acc. under press.", 0.60),
        ("Long balls per match",   0.40),
    ],
    key_evidence=[
        "Lane-level bypass share vs league average.",
        "First-pass type that opens each lane.",
    ],
    also_shown=[
        "Which pressing unit is beaten per lane.",
        "Build-up structure (3\u2011back vs 4\u2011back) split.",
    ],
)


# ═════════════════════════════════════════════════════════════════
# SLIDE 05 — Q4: Post-break shot zones
# ═════════════════════════════════════════════════════════════════
build_question_slide(
    num=5,
    q_tag="Question 04",
    title="Which zones do opponents find shots from after breaking the press?",
    subtitle="This question defines the genuine cost: post-break shot zones and their danger level.",
    user_asks="Once the press is broken, from where do opponents find shots and are those shots dangerous?",
    analyst_answer="Name the post-break shot cluster (edge-of-box central / half-space pull-back / far-post cross) and whether xG per shot is high or low.",
    decision_value="Quantifies the real cost of a broken press and informs rest-defence positioning and pressing-line risk-reward.",
    stat_cards=[
        ("xG per broken press", "0.11"),
        ("Shots per broken press", "0.34"),
        ("Big-chance share", "22%"),
        ("Goals after broken / match", "0.41"),
    ],
    main_caption="Shot map following broken-press sequences (bypass arrow \u2192 shot)",
    pitch_dots=[(0.78, 0.40), (0.82, 0.55), (0.75, 0.65),
                (0.85, 0.50), (0.80, 0.48)],
    supporting_title="Supporting visuals",
    supporting_rows=[
        ("xG per shot",          0.75),
        ("Time to shot",          0.55),
        ("Open-play share",       0.80),
    ],
    key_evidence=[
        "Heatmap of shot origin after broken press.",
        "League comparison of post-break xG per broken press.",
    ],
    also_shown=[
        "Which defensive unit broke first (high line / DM / FB).",
        "Shot type split: open play / cross / through-ball.",
    ],
)


# ── Save ─────────────────────────────────────────────────────────
OUT = "Pressing_Analyst_Storyboard.pptx"
prs.save(OUT)
print(f"Done! {OUT} \u2014 {len(prs.slides)} slides")
