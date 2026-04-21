"""
Build the corrected 14-page pressing analyst presentation.
Matches the structure of pressingfinalr.pdf with fixes applied.
Run:  python build_final_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette (dark green theme matching PDF) ────────────────────
BG      = RGBColor(0, 44, 28)       # deep green background
ORANGE  = RGBColor(255, 100, 0)     # heading orange
WHITE   = RGBColor(255, 255, 255)
BODY    = RGBColor(210, 210, 210)   # body text
MUTED   = RGBColor(160, 160, 160)
GREEN   = RGBColor(0, 200, 120)     # accent green
RED_ACC = RGBColor(255, 90, 90)
FONT    = "Calibri"
MONO    = "Consolas"

# ── Slide dimensions (16:9) ────────────────────────────────────
SW, SH = Inches(13.333), Inches(7.5)
M = Inches(0.6)

# ── Helpers ────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
blank = prs.slide_layouts[6]  # blank layout


def setup(slide):
    """Set dark green background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def tb(slide, left, top, width, height, text,
       sz=18, color=WHITE, bold=False, font=FONT, align=PP_ALIGN.LEFT):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox


def bullet_slide(slide, title, subtitle, bullets, title_color=ORANGE,
                 bullet_size=14, start_y=None):
    """Title + subtitle + bullet list."""
    tb(slide, M, Inches(0.4), Inches(12), Inches(0.6),
       title, sz=28, color=title_color, bold=True)
    if subtitle:
        tb(slide, M, Inches(1.0), Inches(12), Inches(0.5),
           subtitle, sz=18, color=WHITE)
    y = start_y or Inches(1.8)
    for b in bullets:
        tb(slide, Inches(1.0), y, Inches(11), Inches(0.4),
           f"\u2022  {b}", sz=bullet_size, color=BODY)
        y += Inches(0.5)


def section_block(slide, left, top, width, text, sz=11, color=BODY, font=FONT):
    """Multi-line text block."""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = font
        p.space_after = Pt(2)
        # Color overrides for section headers
        stripped = line.strip()
        if stripped == "Strengths:":
            p.font.color.rgb = GREEN
            p.font.bold = True
        elif stripped == "Neutral:":
            p.font.color.rgb = WHITE
            p.font.bold = True
        elif stripped == "Weaknesses:":
            p.font.color.rgb = RED_ACC
            p.font.bold = True
        elif stripped.startswith("Context:"):
            p.font.color.rgb = GREEN
    return txBox


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); setup(s)
tb(s, M, Inches(2.5), Inches(12), Inches(1.0),
   "Context Engineering", sz=40, color=WHITE)
tb(s, M, Inches(3.3), Inches(12), Inches(0.8),
   "Pressing Analyst", sz=40, color=ORANGE, bold=True)
tb(s, M, Inches(4.8), Inches(12), Inches(0.4),
   "Kaan G\u00fcler, Hugo Vicente, Sara Bentelli", sz=14, color=MUTED)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — START WITH FOOTBALL: WHAT DO WE MEASURE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); setup(s)
bullet_slide(s, "Start with football",
    "What do we want to measure?", [
        "How aggressively does the team press in the opponent\u2019s half?",
        "How often do they win the ball back through pressing?",
        "Do they force the opponent into long, risky passes?",
        "How often does the opponent escape the press?",
        "What happens when the press is broken \u2014 does it lead to goals against?",
        "Can the opponent still pass cleanly or break lines under pressure?",
    ])


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — START WITH FOOTBALL: WALL OR GAMBLE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); setup(s)
tb(s, M, Inches(0.4), Inches(12), Inches(0.6),
   "Start with football", sz=28, color=ORANGE, bold=True)
tb(s, M, Inches(1.2), Inches(12), Inches(0.6),
   "Is our pressing a wall or a gamble?", sz=24, color=WHITE)
tb(s, M, Inches(2.2), Inches(11), Inches(2.0),
   "A strong press should do more than apply pressure. It should disrupt build-up, "
   "force rushed decisions, and win the ball high without leaving the team exposed. "
   "This project evaluates that balance by measuring both the benefits of pressing "
   "and the risks when opponents escape it.",
   sz=14, color=BODY)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — DATA TO TEXT (FIXED: metric names, no "nn")
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); setup(s)
tb(s, M, Inches(0.4), Inches(12), Inches(0.6),
   "Data To Text", sz=28, color=ORANGE, bold=True)

# Table with corrected metric names
table_data = [
    ("Metric", "Z Score \u2192 Level", "Consequence Sentence"),
    ("Pressing Chains Per Match",
     "Z > 1.5 \u2192 Outstanding",
     "The team starts pressing sequences in the opposition half far more than any other side."),
    ("Regains In Opp. Half Per Match",
     "1.0 < Z < 1.5 \u2192 Excellent",
     "The team wins the ball back in the opposition half significantly more often than most."),
    ("Opp. Long Balls Forced Per Match",
     "0.5 < Z < 1.0 \u2192 Good",
     "The press forces opponents into long balls more often than average."),
    ("Press Break Rate In Opp. Half",
     "\u22120.5 < Z < 0.5 \u2192 Average",
     "Opponents escape the press in the opposition half at a typical rate."),
    ("Opp. Goals After Broken Press Per Match",
     "\u22121.0 < Z < \u22120.5 \u2192 Below Average",
     "When the press is broken, opponents score more often than against most teams."),
    ("PPDA",
     "1.0 < Z < 1.5 \u2192 Excellent",
     "The team steps in early, allowing opponents very few passes before a defensive action."),
    ("Opp. Pass Accuracy Under Pressure",
     "Z < \u22121.0 \u2192 Poor",
     "Opponents pass comfortably under this team\u2019s pressure \u2014 the press does little to disrupt passing accuracy."),
    ("Opp. Line-Breaking Pass Rate Under Pressure",
     "0.5 < Z < 1.0 \u2192 Good",
     "Opponents break through the defensive lines under this team\u2019s pressure less often than against most sides."),
]

rows, cols = len(table_data), 3
tbl = s.shapes.add_table(rows, cols, M, Inches(1.2), Inches(12.1), Inches(5.5)).table
tbl.columns[0].width = Inches(3.5)
tbl.columns[1].width = Inches(3.0)
tbl.columns[2].width = Inches(5.6)

for r_idx, row_data in enumerate(table_data):
    for c_idx, cell_text in enumerate(row_data):
        cell = tbl.cell(r_idx, c_idx)
        cell.text = cell_text
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10 if r_idx > 0 else 11)
        p.font.name = FONT
        p.font.color.rgb = WHITE if r_idx == 0 else BODY
        p.font.bold = (r_idx == 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 60, 40) if r_idx == 0 else RGBColor(0, 50, 35) if r_idx % 2 == 1 else RGBColor(0, 44, 28)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — SYSTEM PROMPT
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); setup(s)
tb(s, M, Inches(0.4), Inches(12), Inches(0.6),
   "System Prompt", sz=28, color=ORANGE, bold=True)

prompt_lines = [
    '"You are a tactical data analyst embedded in a Premier League coaching staff."',
    '"You write concise pressing briefings for the head coach."',
    '"Your language is direct and action-oriented \u2014 you focus on what the data means tactically."',
    '"You write in British English and refer to the sport as football."',
]
y = Inches(1.5)
for line in prompt_lines:
    tb(s, Inches(1.0), y, Inches(11), Inches(0.4), line, sz=13, color=WHITE, font=MONO)
    y += Inches(0.45)

tb(s, M, Inches(4.5), Inches(11), Inches(1.5),
   "The purpose of this prompt is to transform complex match data into concise, "
   "action-oriented tactical briefings for immediate on-pitch application at "
   "Premier League standards.",
   sz=13, color=MUTED)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — Q&A PAIRS EXAMPLES
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); setup(s)
tb(s, M, Inches(0.4), Inches(12), Inches(0.6),
   "Q&A Pairs Examples", sz=28, color=ORANGE, bold=True)

qa_examples = [
    ("General",
     "What is team pressing in this model?",
     "In this model, team pressing is about how often a side starts pressure sequences "
     "high up the pitch, how much they disrupt the opponent\u2019s build-up, and what "
     "happens if the opponent escapes. The focus is on the opposition half and on what "
     "the opponent is still able to do under pressure."),
    ("Chains per Match",
     "Explain Pressing Chains Per Match and how to interpret it in a pressing analysis.",
     "Pressing Chains Per Match measures how often a team starts pressing sequences in "
     "the opposition half. Higher is better. It is mainly a pressing volume metric rather "
     "than a pure quality metric. High values fit teams that engage opponents high up the "
     "pitch often; low values suggest a more conservative approach or less frequent high pressing."),
    ("high_medium_block_pct",
     "Explain the high/medium block percentage and what it tells us about a team\u2019s pressing.",
     "The high/medium block percentage measures the share of a team\u2019s total out-of-possession "
     "time spent defending in a high or medium block rather than a low block. Neither higher nor "
     "lower is inherently better \u2014 it is a tactical preference, not a quality metric."),
]

y = Inches(1.2)
for cat, q, a in qa_examples:
    tb(s, Inches(0.8), y, Inches(11), Inches(0.3),
       f"Category: {cat}", sz=10, color=ORANGE, bold=True)
    y += Inches(0.3)
    tb(s, Inches(0.8), y, Inches(11), Inches(0.3),
       f"Q: {q}", sz=11, color=WHITE, bold=True)
    y += Inches(0.35)
    tb(s, Inches(0.8), y, Inches(11.2), Inches(1.2),
       f"A: {a}", sz=10, color=BODY)
    y += Inches(1.3)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — SYNTHESIZE (Liverpool example)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); setup(s)
tb(s, M, Inches(0.3), Inches(12), Inches(0.6),
   "Synthesize", sz=28, color=ORANGE, bold=True, font="Calibri")

synth_text = """Strengths:
\u2022 Opp. Line-Breaking Pass Rate Under Pressure \u2014 15.58% \u2014 ranked 4th of 20. Opponents rarely play through the defensive lines under this team\u2019s pressure \u2014 the pressing structure holds its shape.
\u2022 PPDA \u2014 7.90 \u2014 ranked 3rd of 20. The team steps in early, allowing opponents very few passes before a defensive action.

Neutral:
\u2022 Pressing Chains Per Match \u2014 19.24 \u2014 ranked 13th of 20. The team starts pressing sequences in the opposition half at a typical rate.
\u2022 Regains In Opp. Half Per Match \u2014 27.87 \u2014 ranked 5th of 20. The team wins the ball back in the opposition half more often than average.
\u2022 Opp. Long Balls Forced Per Match \u2014 8.00 \u2014 ranked 11th of 20. The press forces opponents into long balls at a typical rate.
\u2022 Press Break Rate In Opp. Half \u2014 61.15% \u2014 ranked 7th of 20. Opponents escape the press in the opposition half less often than against most teams.
\u2022 Opp. Pass Accuracy Under Pressure \u2014 88.90% \u2014 ranked 6th of 20. Opponents find it harder to complete passes under this team\u2019s pressure than against most sides.

Weaknesses:
\u2022 Opp. Goals After Broken Press Per Match \u2014 0.21 \u2014 ranked 16th of 20. When the press is broken, opponents frequently score \u2014 the rest-defence is exposed.

Context: Liverpool spent 91.4% of their out-of-possession time defending in a high or medium block, ranking 10th out of 20 teams in the league. This is not a quality judgement \u2014 it reflects the team\u2019s defensive positioning preference."""

section_block(s, Inches(0.8), Inches(1.0), Inches(11.5), synth_text, sz=10)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — FEW SHOT EXAMPLE (FIXED Liverpool output)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); setup(s)
tb(s, M, Inches(0.4), Inches(12), Inches(0.6),
   "Few Shot Example", sz=28, color=ORANGE, bold=True)

# Fixed output: Strengths in sentence 2, Weaknesses in sentence 3
fixed_output = (
    "Liverpool press at a standard frequency but from a high defensive line, "
    "choosing when to engage rather than sustaining constant pressure. "
    "The press stands out for early intervention and structural discipline \u2014 "
    "opponents are allowed very few passes before a defensive action arrives, "
    "and rarely play through the pressing lines when engaged. "
    "When the press is beaten, opponents score more often than against most teams, "
    "exposing a rest-defence that struggles to contain the damage behind the initial "
    "engagement. "
    "Overall, Liverpool back the quality of their front-foot pressure to compensate "
    "for the vulnerability it creates when bypassed."
)

tb(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.5),
   fixed_output, sz=14, color=BODY)

# Annotation showing sentence mapping
annotations = [
    ("Sentence 1: Pressing identity", "Neutral \u2192 standard frequency, high line"),
    ("Sentence 2: Strengths", "Strengths \u2192 early intervention (PPDA 3rd), structural discipline (Line-Breaking 4th)"),
    ("Sentence 3: Weaknesses", "Weaknesses \u2192 goals after broken press (16th of 20)"),
    ("Sentence 4: Trade-off", "Overall summary of what the approach prioritises and risks"),
]
y = Inches(5.0)
for label, detail in annotations:
    tb(s, Inches(0.8), y, Inches(3), Inches(0.3), label, sz=9, color=GREEN, bold=True)
    tb(s, Inches(3.8), y, Inches(8.5), Inches(0.3), detail, sz=9, color=MUTED)
    y += Inches(0.35)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — NEW: PROMPT INSTRUCTION
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); setup(s)
tb(s, M, Inches(0.4), Inches(12), Inches(0.6),
   "Prompt Instruction", sz=28, color=ORANGE, bold=True)
tb(s, M, Inches(1.0), Inches(12), Inches(0.4),
   "The format instruction sent to the LLM before each team profile", sz=13, color=MUTED)

instruction_text = (
    "Please use the pressing profile enclosed with ``` to give a concise "
    "4-sentence briefing of this team\u2019s pressing style.\n\n"
    "The profile is organised into Strengths, Neutral, and Weaknesses sections.\n\n"
    "Sentence 1: describe the team\u2019s pressing identity \u2014 how and where they press.\n"
    "Sentence 2: describe what the press does well, using the Strengths section.\n"
    "Sentence 3: describe where the press is limited or vulnerable, using the Weaknesses section. "
    "If there are no weaknesses, describe which areas are merely average rather than inventing problems.\n"
    "Sentence 4: summarise the overall trade-off this pressing approach makes.\n\n"
    "Rules:\n"
    "\u2022 Do not invent consequences not supported by the data.\n"
    "\u2022 Do not use metric names, level labels (outstanding, excellent, good, average, below average, poor), "
    "or parenthetical references.\n"
    "\u2022 Write as a tactical analyst briefing coaching staff."
)

section_block(s, Inches(0.8), Inches(1.6), Inches(11.5), instruction_text, sz=12, color=WHITE)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — NEW: PIPELINE OVERVIEW
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); setup(s)
tb(s, M, Inches(0.4), Inches(12), Inches(0.6),
   "Pipeline", sz=28, color=ORANGE, bold=True)
tb(s, M, Inches(1.0), Inches(12), Inches(0.4),
   "How raw metrics become a tactical briefing", sz=14, color=MUTED)

steps = [
    ("\u2776", "Raw Metrics",
     "8 pressing metrics from match data (chains, regains, PPDA, break rate, etc.)"),
    ("\u2777", "Z-Score + Invert",
     "Standardise vs league. Invert 5 negative metrics so positive z always = better."),
    ("\u2778", "Group by Threshold",
     "z > 1.0 \u2192 Strength  |  \u22120.5 to 1.0 \u2192 Neutral  |  z < \u22120.5 \u2192 Weakness"),
    ("\u2779", "Consequence Sentences",
     "48 pre-written sentences (8 metrics \u00d7 6 levels). LLM never interprets z-scores directly."),
    ("\u277a", "Attach Value + Rank",
     "Concrete numbers ground the LLM: \u201c60.53% \u2014 ranked 4th of 20\u201d"),
    ("\u277b", "Structured Profile",
     "Strengths / Neutral / Weaknesses sections + block height context"),
    ("\u277c", "Prompt Stack",
     "System prompt + 19 Q&A pairs + prompt instruction + 4 few-shot examples + team profile"),
    ("\u277d", "LLM Output",
     "4-sentence tactical briefing. Strengths \u2192 sentence 2, Weaknesses \u2192 sentence 3."),
]

y = Inches(1.7)
for icon, title, detail in steps:
    tb(s, Inches(0.8), y, Inches(0.5), Inches(0.4), icon, sz=16, color=GREEN, bold=True)
    tb(s, Inches(1.4), y, Inches(2.5), Inches(0.4), title, sz=13, color=WHITE, bold=True)
    tb(s, Inches(4.0), y, Inches(8.5), Inches(0.4), detail, sz=11, color=BODY)
    y += Inches(0.58)


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — BENCHMARK: FULHAM
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); setup(s)
tb(s, M, Inches(0.4), Inches(12), Inches(0.6),
   "Benchmark", sz=20, color=ORANGE, bold=True)
tb(s, Inches(2.5), Inches(0.4), Inches(5), Inches(0.6),
   " \u2014 Fulham", sz=20, color=WHITE)
tb(s, M, Inches(1.5), Inches(12), Inches(4.5),
   "[Insert Fulham strip chart screenshot here]", sz=14, color=MUTED,
   align=PP_ALIGN.CENTER)
tb(s, M, Inches(5.5), Inches(12), Inches(1.5),
   "[Insert Fulham generated description here]", sz=11, color=MUTED,
   align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDES 12-15 — EXAMPLES (Bournemouth x2, Ipswich, Arsenal)
# ════════════════════════════════════════════════════════════════
example_teams = [
    "Bournemouth", "Bournemouth", "Ipswich Town", "Arsenal"
]
for team in example_teams:
    s = prs.slides.add_slide(blank); setup(s)
    tb(s, M, Inches(0.4), Inches(12), Inches(0.6),
       "Examples", sz=20, color=ORANGE, bold=True)
    tb(s, Inches(2.2), Inches(0.4), Inches(5), Inches(0.6),
       f" \u2014 {team}", sz=20, color=WHITE)
    tb(s, M, Inches(1.5), Inches(12), Inches(4.5),
       f"[Insert {team} strip chart screenshot here]", sz=14, color=MUTED,
       align=PP_ALIGN.CENTER)
    tb(s, M, Inches(5.5), Inches(12), Inches(1.5),
       f"[Insert {team} generated description here]", sz=11, color=MUTED,
       align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 16 — CLOSING
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); setup(s)
tb(s, M, Inches(3.0), Inches(12), Inches(1.0),
   "Ideas? Comments? Questions?", sz=32, color=WHITE, align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────
out = "Pressing_Analyst_final.pptx"
prs.save(out)
print(f"Saved {out} ({len(prs.slides)} slides)")
